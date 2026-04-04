"""
統計学1回_v2.pptx 生成スクリプト
標準本文24pt、全スライドをテキスト/図形で再構築 + シラバス2.0版 2スライド追加
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY    = RGBColor(0x1F, 0x38, 0x64)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x44, 0x72, 0xC4)
GREEN   = RGBColor(0x70, 0xAD, 0x47)
ORANGE  = RGBColor(0xFF, 0xC0, 0x00)
RED     = RGBColor(0xC0, 0x00, 0x00)
PURPLE  = RGBColor(0x70, 0x30, 0xA0)
DKGRAY  = RGBColor(0x26, 0x26, 0x26)
MGRAY   = RGBColor(0x59, 0x59, 0x59)
LGRAY   = RGBColor(0xF2, 0xF2, 0xF2)
LBLUE   = RGBColor(0xBD, 0xD7, 0xEE)
LCYAN   = RGBColor(0xE1, 0xF5, 0xFE)
LORANGE = RGBColor(0xFF, 0xF2, 0xCC)
LGREEN  = RGBColor(0xE2, 0xEF, 0xDA)
LRED    = RGBColor(0xFF, 0xD0, 0xD0)
TEAL    = RGBColor(0x00, 0x70, 0x6E)
LTEAL   = RGBColor(0xD0, 0xF0, 0xEE)
GOLD    = RGBColor(0xB8, 0x86, 0x00)

W = Inches(13.333); H = Inches(7.5)
FONT = "メイリオ"; STD = 24

prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(sl, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    shp = sl.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line: shp.line.color.rgb = line; shp.line.width = lw
    else: shp.line.fill.background()
    return shp

def txb(sl, l, t, w, h, lines, sz=STD, bold=False, color=DKGRAY,
        align=PP_ALIGN.LEFT, wrap=True):
    box = sl.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    if isinstance(lines, str): lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(2)
        run = p.add_run()
        run.text = line; run.font.size = Pt(sz)
        run.font.bold = bold; run.font.color.rgb = color
        run.font.name = FONT
    return box

def header(sl, title, sub=None):
    rect(sl, 0, 0, W, Inches(0.85), fill=NAVY)
    txb(sl, Inches(0.35), Inches(0.08), W-Inches(0.7), Inches(0.72),
        title, sz=26, bold=True, color=WHITE)
    if sub:
        txb(sl, Inches(0.35), Inches(0.92), W-Inches(0.7), Inches(0.5),
            sub, sz=18, bold=True, color=BLUE)

def th(sl, y, cols, rh=Inches(0.55)):
    x = Inches(0.4)
    for txt, w in cols:
        rect(sl, x, y, w, rh, fill=NAVY, line=WHITE, lw=Pt(0.5))
        txb(sl, x+Inches(0.05), y, w-Inches(0.1), rh,
            [txt], sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w

def tr(sl, y, cells, rh=Inches(0.65), even=True):
    bg = LGRAY if even else WHITE; x = Inches(0.4)
    for txt, w, *rest in cells:
        clr = rest[0] if rest else bg
        rect(sl, x, y, w, rh, fill=clr, line=BLUE, lw=Pt(0.5))
        txb(sl, x+Inches(0.05), y, w-Inches(0.1), rh,
            [txt], sz=18, color=DKGRAY, align=PP_ALIGN.CENTER)
        x += w

# ════ S1: タイトル ════
s = slide()
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.06), fill=BLUE)
txb(s, Inches(1.5), Inches(0.7), Inches(10.3), Inches(1.3),
    "統計学", sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0),
    "第 1 回：統計の基本", sz=32, color=ORANGE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.6),
    "齋藤 邦彦", sz=22, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.5),
    "データサイエンスの基礎となる統計学 — 平均・分散・Σ記号から始める", sz=18, color=LBLUE, align=PP_ALIGN.CENTER)

# ════ S2: 統計とデータサイエンス ════
s = slide()
header(s, "統計とデータサイエンス", "Statistics & Data Science")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(1.0), fill=NAVY)
txb(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5),
    "データサイエンスにおける統計学の位置付け", sz=22, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.4),
    "金融・経営・財務会計・HRなどのビジネスデータを科学的に分析するための基礎学問", sz=19, color=LBLUE)

concepts = [
    ("データサイエンス", BLUE,
     ["データ処理・分析の技術体系", "ビジネス課題の解決手法", "AI・機械学習の基盤"]),
    ("統計学", GREEN,
     ["データの記述・要約（基本統計量）", "推測・仮説検定・相関分析", "確率論・分布理論"]),
    ("プログラミング", ORANGE,
     ["Excel（関数・ピボット・グラフ）", "Python（pandas・numpy・scipy）", "データ可視化・自動化"]),
]
for i, (title, c, items) in enumerate(concepts):
    x = Inches(0.4) + Inches(4.2)*i
    rect(s, x, Inches(2.15), Inches(4.0), Inches(0.65), fill=c)
    txb(s, x+Inches(0.1), Inches(2.2), Inches(3.8), Inches(0.6),
        title, sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.8), Inches(4.0), Inches(2.0), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.15), Inches(2.9), Inches(3.7), Inches(1.85),
        items, sz=20, color=DKGRAY)

txb(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.45),
    "第1-7回のカリキュラム概要：", sz=20, bold=True, color=NAVY)
rows_prog = [
    ("第1回", "統計の基本", "平均・最大・最小・範囲・Σ記号"),
    ("第2回", "ばらつきの統計量", "分散・標準偏差・変動係数"),
    ("第3-4回", "確率・分布", "正規分布・二項分布・t分布"),
    ("第5-6回", "推測統計", "区間推定・仮説検定"),
    ("第7回", "相関・回帰", "相関係数・単回帰分析"),
]
th(s, Inches(5.55), [("回", Inches(1.5)), ("テーマ", Inches(3.0)), ("主な内容", Inches(7.7))], rh=Inches(0.48))
for i, (r, t, c) in enumerate(rows_prog):
    tr(s, Inches(6.03)+Inches(0.52)*i, [(r, Inches(1.5)), (t, Inches(3.0)), (c, Inches(7.7))],
       rh=Inches(0.5), even=(i%2==0))

# ════ S3: Σ記号 ════
s = slide()
header(s, "合計を求める記号：Σ（シグマ）", "Summation Notation")
# 大きなΣ
rect(s, Inches(0.4), Inches(1.05), Inches(5.8), Inches(5.7), fill=LCYAN, line=BLUE, lw=Pt(2))
txb(s, Inches(0.6), Inches(1.2), Inches(5.4), Inches(0.6),
    "Σ（シグマ）の意味", sz=22, bold=True, color=NAVY)
txb(s, Inches(1.0), Inches(1.9), Inches(4.8), Inches(1.4),
    "Σ", sz=100, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txb(s, Inches(0.6), Inches(3.4), Inches(5.4), Inches(0.5),
    "読み方：シグマ（Sigma）/ サメーション（Summation）", sz=18, color=DKGRAY)
txb(s, Inches(0.6), Inches(4.0), Inches(5.4), Inches(1.5),
    ["意味：添え字の初めの値から",
     "　　　データの件数までの数を",
     "　　　計算式に代入して合計する"],
    sz=20, color=DKGRAY)

rect(s, Inches(6.5), Inches(1.05), Inches(6.4), Inches(5.7), fill=LGRAY, line=NAVY, lw=Pt(1.5))
txb(s, Inches(6.7), Inches(1.15), Inches(6.0), Inches(0.6),
    "Σ の構成要素", sz=22, bold=True, color=NAVY)
parts = [
    ("上限（n）", "データの件数（最後の添え字値）", BLUE),
    ("下限（i=1）", "添え字の初めの値", GREEN),
    ("計算式（xi）", "各データ要素に対する演算", ORANGE),
]
for i, (part, desc, c) in enumerate(parts):
    y = Inches(1.95) + Inches(1.3)*i
    rect(s, Inches(6.6), y, Inches(2.4), Inches(1.1), fill=c)
    txb(s, Inches(6.7), y+Inches(0.15), Inches(2.2), Inches(0.85),
        part, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(9.0), y, Inches(3.7), Inches(1.1), fill=WHITE, line=c, lw=Pt(1))
    txb(s, Inches(9.1), y+Inches(0.15), Inches(3.5), Inches(0.85),
        desc, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.5), fill=NAVY)
txb(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45),
    "例：Σxi（i=1からn）= x1 + x2 + x3 + … + xn　← 全データの合計",
    sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════ S4: データの件数・合計・範囲 ════
s = slide()
header(s, "基本統計量：件数・合計・範囲", "Basic Statistics: Count, Sum, Range")
defs = [
    ("データの件数 N", "データの総数（サンプルサイズ）", "N = 7（7つのデータがある場合）", BLUE),
    ("データの合計 Σxi", "すべての要素を足し合わせた値", "x=[6,8,5,9,8,4,7] → Σxi = 47", GREEN),
    ("データの範囲 R", "データの最大値と最小値の差", "R = xmax − xmin = 9 − 4 = 5", ORANGE),
]
for i, (name, defn, ex, c) in enumerate(defs):
    y = Inches(1.1) + Inches(1.7)*i
    rect(s, Inches(0.4), y, Inches(3.5), Inches(1.55), fill=c)
    txb(s, Inches(0.5), y+Inches(0.15), Inches(3.3), Inches(0.6),
        name, sz=20, bold=True, color=WHITE)
    rect(s, Inches(3.9), y, Inches(5.5), Inches(1.55), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, Inches(4.0), y+Inches(0.15), Inches(5.3), Inches(0.6),
        defn, sz=20, color=DKGRAY)
    rect(s, Inches(9.4), y, Inches(3.5), Inches(1.55), fill=LCYAN, line=c, lw=Pt(1))
    txb(s, Inches(9.5), y+Inches(0.15), Inches(3.3), Inches(0.6),
        "例：", sz=17, bold=True, color=NAVY)
    txb(s, Inches(9.5), y+Inches(0.5), Inches(3.3), Inches(0.8),
        ex, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.55), fill=LORANGE, line=ORANGE, lw=Pt(1))
txb(s, Inches(0.5), Inches(6.28), Inches(12.3), Inches(0.45),
    "データ x = [ 6, 8, 5, 9, 8, 4, 7 ]　→　N=7, Σxi=47, xmax=9, xmin=4, 範囲R=5",
    sz=20, bold=True, color=GOLD)

# ════ S5: 平均値の定義 ════
s = slide()
header(s, "平均値（算術平均）", "Arithmetic Mean")
# 式の表示
rect(s, Inches(0.4), Inches(1.05), Inches(7.0), Inches(3.5), fill=LCYAN, line=BLUE, lw=Pt(2))
txb(s, Inches(0.6), Inches(1.15), Inches(6.6), Inches(0.6),
    "平均値の定義式", sz=22, bold=True, color=NAVY)

txb(s, Inches(0.6), Inches(1.85), Inches(6.6), Inches(0.55),
    "E = (1/n) × Σxi", sz=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txb(s, Inches(0.6), Inches(2.5), Inches(6.6), Inches(0.4),
    "または", sz=18, color=MGRAY, align=PP_ALIGN.CENTER)
txb(s, Inches(0.6), Inches(2.95), Inches(6.6), Inches(0.55),
    "E = (x1 + x2 + … + xn) / n", sz=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# 変数の説明
txb(s, Inches(0.6), Inches(3.65), Inches(6.6), Inches(0.45),
    "記号の説明：", sz=18, bold=True, color=NAVY)
vars_desc = [("E", "平均値（Expected Value / Mean）"),
             ("n", "データの件数"),
             ("xi", "個々のデータの値（i番目）")]
for i, (v, d) in enumerate(vars_desc):
    rect(s, Inches(0.7), Inches(4.15)+Inches(0.6)*i, Inches(0.6), Inches(0.52), fill=NAVY)
    txb(s, Inches(0.73), Inches(4.2)+Inches(0.6)*i, Inches(0.55), Inches(0.5),
        v, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, Inches(1.4), Inches(4.2)+Inches(0.6)*i, Inches(5.5), Inches(0.5),
        d, sz=19, color=DKGRAY)

# 右側：特徴
rect(s, Inches(7.6), Inches(1.05), Inches(5.3), Inches(3.5), fill=LGRAY, line=NAVY, lw=Pt(1.5))
txb(s, Inches(7.8), Inches(1.15), Inches(5.0), Inches(0.6),
    "平均値の特徴", sz=22, bold=True, color=NAVY)
features = [
    ("最も代表的な代表値", "データの中心的傾向を示す"),
    ("全データの影響を受ける", "外れ値に影響されやすい"),
    ("「おもに算術平均」を指す", "他に幾何平均・調和平均あり"),
    ("期待値とも呼ばれる", "確率論・統計学の基礎"),
]
for i, (f, d) in enumerate(features):
    y = Inches(1.85) + Inches(0.65)*i
    rect(s, Inches(7.7), y, Inches(5.1), Inches(0.55), fill=LBLUE if i%2==0 else WHITE, line=BLUE, lw=Pt(0.3))
    txb(s, Inches(7.8), y+Inches(0.05), Inches(4.9), Inches(0.5),
        f"▶ {f}：{d}", sz=18, color=DKGRAY)

txb(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.45),
    "他の代表値との比較：", sz=20, bold=True, color=NAVY)
th(s, Inches(5.55), [("代表値", Inches(2.5)), ("定義", Inches(4.5)), ("特徴", Inches(5.2))], rh=Inches(0.5))
rep = [
    ("算術平均", "全データの和 ÷ 件数", "最も一般的。外れ値に敏感"),
    ("中央値", "データを並べた中央の値", "外れ値に頑健。順序データ向き"),
    ("最頻値", "最も頻繁に現れる値", "カテゴリデータ向き"),
]
for i, (a, b, c) in enumerate(rep):
    tr(s, Inches(6.05)+Inches(0.55)*i, [(a, Inches(2.5)), (b, Inches(4.5)), (c, Inches(5.2))],
       rh=Inches(0.52), even=(i%2==0))

# ════ S6: 平均値の計算例 ════
s = slide()
header(s, "平均値の計算例", "Calculating the Mean")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=BLUE)
txb(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6),
    "数列データ x = [ 6, 8, 5, 9, 8, 4, 7 ] の基本統計量を求める",
    sz=21, bold=True, color=WHITE)

# データ可視化（棒グラフ風）
vals = [6, 8, 5, 9, 8, 4, 7]
bar_w = Inches(1.2)
bar_base = Inches(4.5)
max_h = Inches(2.8)
colors_bar = [BLUE, BLUE, BLUE, RED, BLUE, GREEN, BLUE]
for i, v in enumerate(vals):
    x = Inches(0.6) + bar_w * i
    bar_h = max_h * v / 10
    y = bar_base - bar_h
    rect(s, x, y, bar_w - Inches(0.1), bar_h, fill=colors_bar[i])
    txb(s, x, y - Inches(0.4), bar_w - Inches(0.1), Inches(0.4),
        str(v), sz=20, bold=True, color=DKGRAY, align=PP_ALIGN.CENTER)
    txb(s, x, bar_base + Inches(0.05), bar_w - Inches(0.1), Inches(0.4),
        f"x{i+1}", sz=17, color=MGRAY, align=PP_ALIGN.CENTER)
# 平均線
rect(s, Inches(0.5), bar_base - max_h * 6.71 / 10, Inches(8.9), Inches(0.04), fill=ORANGE)
txb(s, Inches(9.5), bar_base - max_h * 6.71 / 10 - Inches(0.2), Inches(3.5), Inches(0.4),
    "← 平均値 6.71", sz=18, bold=True, color=ORANGE)

# 計算過程
rect(s, Inches(9.8), Inches(1.05), Inches(3.1), Inches(5.5), fill=LGRAY, line=NAVY, lw=Pt(1.5))
txb(s, Inches(9.9), Inches(1.15), Inches(2.9), Inches(0.5),
    "計算過程", sz=20, bold=True, color=NAVY)
calc = [
    "① 合計を求める",
    "6+8+5+9+8+4+7 = 47",
    "",
    "② 件数で割る",
    "47 ÷ 7 = 6.714…",
    "",
    "③ 結果",
    "平均値 E ≈ 6.71",
]
txb(s, Inches(9.9), Inches(1.75), Inches(2.9), Inches(4.0),
    calc, sz=19, color=DKGRAY)

# 基本統計量まとめ
stats_list = [
    ("平均値 E", "6.71", BLUE),
    ("最大値 xmax", "9", RED),
    ("最小値 xmin", "4", GREEN),
    ("範囲 R", "9 − 4 = 5", ORANGE),
    ("件数 N", "7", PURPLE),
]
txb(s, Inches(0.4), Inches(4.65), Inches(9.2), Inches(0.5),
    "基本統計量まとめ：", sz=19, bold=True, color=NAVY)
for i, (stat, val, c) in enumerate(stats_list):
    x = Inches(0.4) + Inches(1.8)*i
    rect(s, x, Inches(5.2), Inches(1.7), Inches(0.6), fill=c)
    txb(s, x+Inches(0.05), Inches(5.25), Inches(1.6), Inches(0.55),
        stat, sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(5.8), Inches(1.7), Inches(0.55), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.05), Inches(5.85), Inches(1.6), Inches(0.5),
        val, sz=20, bold=True, color=c, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.65), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6),
    "Excel関数：=AVERAGE(範囲)　=MAX(範囲)　=MIN(範囲)　=COUNT(範囲)　=MAX()-MIN()",
    sz=19, bold=True, color=GOLD)

# ════ S7: 幾何平均・調和平均 ════
s = slide()
header(s, "幾何平均と調和平均", "Geometric & Harmonic Mean")
rect(s, Inches(0.4), Inches(1.05), Inches(6.0), Inches(5.8), fill=LCYAN, line=BLUE, lw=Pt(2))
txb(s, Inches(0.6), Inches(1.15), Inches(5.6), Inches(0.6),
    "幾何平均（Geometric Mean）", sz=22, bold=True, color=BLUE)
txb(s, Inches(0.6), Inches(1.85), Inches(5.6), Inches(0.55),
    "G = (X1 × X2 × … × Xn)^(1/n)", sz=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txb(s, Inches(0.6), Inches(2.5), Inches(5.6), Inches(0.45),
    "Excel：=GEOMEAN(範囲)", sz=18, color=MGRAY, align=PP_ALIGN.CENTER)
txb(s, Inches(0.6), Inches(3.05), Inches(5.6), Inches(0.5),
    "■ 使いどき：", sz=19, bold=True, color=NAVY)
txb(s, Inches(0.6), Inches(3.6), Inches(5.6), Inches(1.5),
    ["・ 比率・成長率・倍率の平均",
     "・ 投資リターンの平均計算",
     "・ 物価指数・人口増加率"],
    sz=19, color=DKGRAY)
txb(s, Inches(0.6), Inches(5.2), Inches(5.6), Inches(0.5),
    "■ 計算例：家賃が1年で20%・10%・15%上昇", sz=17, bold=True, color=NAVY)
txb(s, Inches(0.6), Inches(5.75), Inches(5.6), Inches(0.9),
    ["G = (1.20 × 1.10 × 1.15)^(1/3) ≈ 1.148",
     "→ 年平均約14.8%上昇"],
    sz=18, color=DKGRAY)

rect(s, Inches(6.8), Inches(1.05), Inches(6.1), Inches(5.8), fill=LGREEN, line=GREEN, lw=Pt(2))
txb(s, Inches(7.0), Inches(1.15), Inches(5.7), Inches(0.6),
    "調和平均（Harmonic Mean）", sz=22, bold=True, color=TEAL)
txb(s, Inches(7.0), Inches(1.85), Inches(5.7), Inches(0.55),
    "H = n / (1/X1 + 1/X2 + … + 1/Xn)", sz=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txb(s, Inches(7.0), Inches(2.5), Inches(5.7), Inches(0.45),
    "Excel：=HARMEAN(範囲)", sz=18, color=MGRAY, align=PP_ALIGN.CENTER)
txb(s, Inches(7.0), Inches(3.05), Inches(5.7), Inches(0.5),
    "■ 使いどき：", sz=19, bold=True, color=TEAL)
txb(s, Inches(7.0), Inches(3.6), Inches(5.7), Inches(1.5),
    ["・ 速度・単価などの比率の平均",
     "・ 「時間あたり」の量を扱う場面",
     "・ 算術平均では誤差が出るとき"],
    sz=19, color=DKGRAY)
txb(s, Inches(7.0), Inches(5.2), Inches(5.7), Inches(0.5),
    "■ 計算例：100kmを時速50・40・60kmで区間走行", sz=17, bold=True, color=TEAL)
txb(s, Inches(7.0), Inches(5.75), Inches(5.7), Inches(0.9),
    ["計算：100 ÷ (20/50 + 50/40 + 30/60) ≈ 46.5km/h",
     "→ 距離加重平均速度 ≈ 46.5 km/h"],
    sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.35), fill=NAVY)
txb(s, Inches(0.5), Inches(6.98), Inches(12.3), Inches(0.3),
    "3つの平均の大小関係：調和平均 H ≦ 幾何平均 G ≦ 算術平均 E（等号は全データが等しい時）",
    sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════ S8: Excel関数 ════
s = slide()
header(s, "Excelで基本統計量を求める", "Excel Functions for Statistics")
txb(s, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.45),
    "ビジネスデータを用いてExcelで基本統計量を計算する方法", sz=20, color=MGRAY)

th(s, Inches(1.6), [("関数名", Inches(3.0)), ("構文", Inches(4.5)), ("用途・説明", Inches(5.2))], rh=Inches(0.52))
excel_fns = [
    ("AVERAGE()", "=AVERAGE(A1:A10)", "算術平均を計算"),
    ("AVERAGEIF()", "=AVERAGEIF(条件範囲,条件,平均範囲)", "条件付き平均（例：女性のみ）"),
    ("MAX()", "=MAX(A1:A10)", "最大値を返す"),
    ("MIN()", "=MIN(A1:A10)", "最小値を返す"),
    ("COUNT()", "=COUNT(A1:A10)", "数値データの件数を数える"),
    ("GEOMEAN()", "=GEOMEAN(A1:A10)", "幾何平均を計算"),
    ("HARMEAN()", "=HARMEAN(A1:A10)", "調和平均を計算"),
    ("STDEV()", "=STDEV(A1:A10)", "標本標準偏差を計算"),
]
for i, (fn, syntax, desc) in enumerate(excel_fns):
    tr(s, Inches(2.12)+Inches(0.56)*i, [(fn, Inches(3.0)), (syntax, Inches(4.5)), (desc, Inches(5.2))],
       rh=Inches(0.54), even=(i%2==0))

rect(s, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.6), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.55),
    "AVERAGEIF の条件の書き方：\"=F\"（女性）/ \">=80\"（80以上） / \"=\"&B4（セル参照）",
    sz=18, bold=True, color=GOLD)

# ════ S9: まとめ ════
s = slide()
header(s, "第1回 まとめ", "Summary")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.55), fill=BLUE)
txb(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
    "今日学んだキーワードと概念", sz=20, bold=True, color=WHITE)

keywords = [
    ("Σ（シグマ）記号", "合計を簡潔に表す数学記号。i=1からnまでの計算式を合計する表現"),
    ("データの件数 N", "標本サイズ。全データの数を表す"),
    ("データの合計 Σxi", "全データ要素の和。Σ記号を用いて表現"),
    ("データの範囲 R", "最大値 − 最小値。データのばらつきの最も単純な指標"),
    ("算術平均 E", "最も一般的な代表値。Σxi ÷ n で計算"),
    ("幾何平均 G", "比率・成長率の平均。(X1×…×Xn)^(1/n)"),
    ("調和平均 H", "速度・単価など比率の逆数の平均。H ≦ G ≦ E"),
]
for i, (kw, desc) in enumerate(keywords):
    y = Inches(1.7) + Inches(0.72)*i
    rect(s, Inches(0.4), y, Inches(3.0), Inches(0.65), fill=NAVY)
    txb(s, Inches(0.5), y+Inches(0.08), Inches(2.8), Inches(0.58),
        kw, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.4), y, Inches(9.5), Inches(0.65), fill=LGRAY if i%2==0 else WHITE, line=BLUE, lw=Pt(0.3))
    txb(s, Inches(3.5), y+Inches(0.08), Inches(9.3), Inches(0.58),
        desc, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.55), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
    "次回：ばらつきの統計量（分散・標準偏差）と正規分布の基礎",
    sz=20, bold=True, color=GOLD)

# ════ S10: 課題 ════
s = slide()
header(s, "課題・演習", "Assignment")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=ORANGE)
txb(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6),
    "【課題】次のクラスデータから平均点を求めよ（初回は提出無し）",
    sz=21, bold=True, color=WHITE)

rect(s, Inches(0.4), Inches(1.85), Inches(12.5), Inches(1.3), fill=LGRAY, line=BLUE, lw=Pt(1))
txb(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
    "問1：あるクラス50人の100点満点の国語テスト結果から平均点を求めよ",
    sz=20, bold=True, color=NAVY)
txb(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.5),
    "ヒント：mean ≈ 59.72, max = 83, min = 41　　（Google Docs・ChatGPT等の利用可）",
    sz=18, color=DKGRAY)

txb(s, Inches(0.4), Inches(3.35), Inches(12.5), Inches(0.5),
    "■ 解き方のステップ：", sz=20, bold=True, color=NAVY)
steps = [
    ("STEP 1", "データを確認する（50個の得点を整理）", BLUE),
    ("STEP 2", "全データの合計 Σxi を求める", GREEN),
    ("STEP 3", "件数 n = 50 で割る（平均 = Σxi ÷ 50）", ORANGE),
    ("STEP 4", "最大値・最小値・範囲も求める", PURPLE),
    ("STEP 5", "Excelで =AVERAGE(範囲) を使って確認する", TEAL),
]
for i, (step, desc, c) in enumerate(steps):
    y = Inches(3.9) + Inches(0.65)*i
    rect(s, Inches(0.5), y, Inches(1.5), Inches(0.58), fill=c)
    txb(s, Inches(0.55), y+Inches(0.08), Inches(1.4), Inches(0.52),
        step, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.0), y, Inches(10.9), Inches(0.58), fill=LGRAY if i%2==0 else WHITE, line=c, lw=Pt(0.5))
    txb(s, Inches(2.1), y+Inches(0.08), Inches(10.7), Inches(0.52),
        desc, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.25), fill=NAVY)
txb(s, Inches(0.5), Inches(7.13), Inches(12.3), Inches(0.22),
    "次回からはレポートとして提出 — 答え合わせは翌週",
    sz=14, color=WHITE, align=PP_ALIGN.CENTER)

# ════ S11: 中央値（Median） ════
s = slide()
header(s, "中央値（メジアン）", "Median")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.9), fill=NAVY)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.45),
    "中央値（Median）= データを昇順に並べたとき、ちょうど中央に位置する値", sz=22, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.3),
    "外れ値（異常に大きい・小さい値）の影響を受けにくく、分布が偏っているデータに有効", sz=17, color=LBLUE)

# 奇数の例
rect(s, Inches(0.4), Inches(2.1), Inches(6.0), Inches(3.8), fill=LCYAN, line=BLUE, lw=Pt(2))
txb(s, Inches(0.6), Inches(2.2), Inches(5.6), Inches(0.5),
    "【奇数個の場合】データ数 n = 7", sz=20, bold=True, color=NAVY)
txb(s, Inches(0.6), Inches(2.8), Inches(5.6), Inches(0.45),
    "データ：4, 5, 6, 7, 8, 9, 100（昇順）", sz=19, color=DKGRAY)
txb(s, Inches(0.6), Inches(3.35), Inches(5.6), Inches(0.45),
    "位置：(7+1)÷2 = 4番目", sz=19, color=DKGRAY)
rect(s, Inches(0.6), Inches(3.9), Inches(5.6), Inches(0.65), fill=BLUE)
txb(s, Inches(0.7), Inches(3.95), Inches(5.4), Inches(0.55),
    "中央値 = 7（4番目の値）", sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(0.6), Inches(4.7), Inches(5.6), Inches(0.5),
    "平均値 = (4+5+6+7+8+9+100)÷7 ≈ 19.9", sz=18, color=RED)
txb(s, Inches(0.6), Inches(5.25), Inches(5.6), Inches(0.5),
    "→ 中央値7が実態をより正確に反映", sz=18, bold=True, color=NAVY)

# 偶数の例
rect(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(3.8), fill=LGREEN, line=GREEN, lw=Pt(2))
txb(s, Inches(7.0), Inches(2.2), Inches(5.6), Inches(0.5),
    "【偶数個の場合】データ数 n = 6", sz=20, bold=True, color=NAVY)
txb(s, Inches(7.0), Inches(2.8), Inches(5.6), Inches(0.45),
    "データ：3, 5, 7, 9, 11, 13（昇順）", sz=19, color=DKGRAY)
txb(s, Inches(7.0), Inches(3.35), Inches(5.6), Inches(0.45),
    "中央の2値：3番目=7、4番目=9", sz=19, color=DKGRAY)
rect(s, Inches(7.0), Inches(3.9), Inches(5.6), Inches(0.65), fill=GREEN)
txb(s, Inches(7.1), Inches(3.95), Inches(5.4), Inches(0.55),
    "中央値 = (7+9)÷2 = 8", sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(7.0), Inches(4.7), Inches(5.6), Inches(0.5),
    "偶数個：中央2値の平均が中央値", sz=18, color=DKGRAY)
txb(s, Inches(7.0), Inches(5.25), Inches(5.6), Inches(0.5),
    "→ 平均値と同じ 8（均一分布の場合）", sz=18, color=MGRAY)

rect(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.55), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.5),
    "Excel関数：=MEDIAN(範囲)　　Python：scores.median()　　計算手順：①昇順並べ替え → ②中央の位置特定",
    sz=19, bold=True, color=GOLD)

# ════ S12: 最頻値（Mode） ════
s = slide()
header(s, "最頻値（モード）", "Mode")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.9), fill=NAVY)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.45),
    "最頻値（Mode）= データの中で最も多く出現する値", sz=22, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.3),
    "カテゴリデータ（性別・血液型・商品種別）や、特定の値が集中するデータに最適", sz=17, color=LBLUE)

rect(s, Inches(0.4), Inches(2.1), Inches(5.8), Inches(2.8), fill=LGRAY, line=PURPLE, lw=Pt(2))
txb(s, Inches(0.6), Inches(2.2), Inches(5.4), Inches(0.5),
    "例1：テスト点数データ", sz=20, bold=True, color=PURPLE)
txb(s, Inches(0.6), Inches(2.75), Inches(5.4), Inches(0.4),
    "70, 80, 80, 80, 90, 90, 95", sz=18, color=DKGRAY)
rect(s, Inches(0.6), Inches(3.25), Inches(5.4), Inches(0.65), fill=PURPLE)
txb(s, Inches(0.7), Inches(3.3), Inches(5.2), Inches(0.55),
    "最頻値 = 80（3回出現）", sz=21, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(0.6), Inches(4.0), Inches(5.4), Inches(0.7),
    ["→ 最も多くの生徒が取った点数",
     "→ 「典型的な成績」を示す"],
    sz=17, color=DKGRAY)

rect(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(2.8), fill=LGRAY, line=ORANGE, lw=Pt(2))
txb(s, Inches(7.0), Inches(2.2), Inches(5.6), Inches(0.5),
    "例2：商品の売れ筋サイズ", sz=20, bold=True, color=ORANGE)
txb(s, Inches(7.0), Inches(2.75), Inches(5.6), Inches(0.4),
    "S, M, M, L, M, XL, L, M（服のサイズ）", sz=17, color=DKGRAY)
rect(s, Inches(7.0), Inches(3.25), Inches(5.6), Inches(0.65), fill=ORANGE)
txb(s, Inches(7.1), Inches(3.3), Inches(5.4), Inches(0.55),
    "最頻値 = M（4回出現）", sz=21, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(7.0), Inches(4.0), Inches(5.6), Inches(0.7),
    ["→ 最も売れている「M」が最頻値",
     "→ 在庫管理・仕入れに活用"],
    sz=17, color=DKGRAY)

txb(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.45),
    "■ 複数の最頻値（双峰・多峰分布）", sz=20, bold=True, color=NAVY)
rect(s, Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.9), fill=LRED, line=RED, lw=Pt(1))
txb(s, Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.35),
    "データ：1, 2, 2, 3, 4, 4, 5 → 最頻値 = 2 と 4（両方3回ではなく2回ずつ）", sz=18, color=DKGRAY)
txb(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.3),
    "→ 最頻値が2つある場合「双峰分布」といい、2つの集団が混在している可能性がある（例：身長の男女別分布）",
    sz=17, color=RED)

rect(s, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.5), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.45),
    "Excel：=MODE(範囲)  または  =MODE.MULT(範囲)（複数最頻値）　　Python：scores.mode()[0]",
    sz=19, bold=True, color=GOLD)

# ════ S13: 平均・中央値・最頻値の比較 ════
s = slide()
header(s, "平均・中央値・最頻値の使い分け", "Mean vs Median vs Mode")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=BLUE)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.5),
    "3つの「代表値」 — データの典型的な値を1つの数で表す統計量", sz=20, bold=True, color=WHITE)

th(s, Inches(1.85), [
    ("", Inches(2.2)),
    ("算術平均 Mean", Inches(3.3)),
    ("中央値 Median", Inches(3.3)),
    ("最頻値 Mode", Inches(3.3)),
], rh=Inches(0.52))
comp_rows = [
    ("計算方法", "全データの合計÷件数", "昇順の中央の値", "最多出現値"),
    ("外れ値の影響", "大きく受ける", "ほとんど受けない", "受けない"),
    ("データ種別", "数値データ", "数値データ", "数値・カテゴリ"),
    ("適したケース", "正規分布・均一分布", "偏った分布・年収", "カテゴリ・売れ筋"),
    ("例：年収データ", "外れ値で高くなる", "実態を正確に反映", "最多の所得層"),
]
for i, row in enumerate(comp_rows):
    y = Inches(2.37) + Inches(0.72)*i
    even = i % 2 == 0
    vals = [(row[0], Inches(2.2)), (row[1], Inches(3.3)), (row[2], Inches(3.3)), (row[3], Inches(3.3))]
    x = Inches(0.4)
    fills = [LGRAY if even else WHITE, LGRAY if even else WHITE,
             LCYAN if even else WHITE, LGREEN if even else WHITE]
    colors = [DKGRAY, DKGRAY, NAVY, GREEN]
    for j, (txt, w) in enumerate(vals):
        rect(s, x, y, w, Inches(0.66), fill=fills[j], line=BLUE, lw=Pt(0.3))
        txb(s, x+Inches(0.05), y+Inches(0.05), w-Inches(0.1), Inches(0.58),
            txt, sz=17, color=colors[j], align=PP_ALIGN.CENTER)
        x += w

rect(s, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.65), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.55),
    "実務での使い分け：年収・地価→中央値 ／ 販売数・成績分析→平均値 ／ 人気商品・血液型→最頻値",
    sz=19, bold=True, color=GOLD)
txb(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5),
    "→ 3つを並べて比較することでデータの分布形状（対称・右偏・左偏）が把握できる",
    sz=18, color=DKGRAY)

# ════ S14: 度数分布表 ════
s = slide()
header(s, "度数分布表とヒストグラム", "Frequency Distribution & Histogram")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=NAVY)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.45),
    "度数分布表：データをいくつかの区間（階級）に分けて、各区間の件数を集計した表",
    sz=20, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.3),
    "大量データの全体像を把握するための基本的な集計方法", sz=17, color=LBLUE)

rect(s, Inches(0.4), Inches(2.0), Inches(5.8), Inches(4.8), fill=LGRAY, line=BLUE, lw=Pt(1.5))
txb(s, Inches(0.6), Inches(2.1), Inches(5.4), Inches(0.45),
    "例：50人のテスト点数（100点満点）", sz=18, bold=True, color=NAVY)
th_cols = [("階級（点数）", Inches(2.2)), ("度数（人）", Inches(1.7)), ("相対度数", Inches(1.7))]
th(s, Inches(2.65), th_cols, rh=Inches(0.48))
freq_rows = [
    ("40〜49点", "3人", "6%"),
    ("50〜59点", "8人", "16%"),
    ("60〜69点", "15人", "30%"),
    ("70〜79点", "14人", "28%"),
    ("80〜89点", "8人", "16%"),
    ("90〜99点", "2人", "4%"),
]
for i, (cls, freq, rel) in enumerate(freq_rows):
    tr(s, Inches(3.13)+Inches(0.52)*i,
       [(cls, Inches(2.2)), (freq, Inches(1.7)), (rel, Inches(1.7))],
       rh=Inches(0.5), even=(i%2==0))

rect(s, Inches(6.5), Inches(2.0), Inches(6.3), Inches(4.8), fill=WHITE, line=BLUE, lw=Pt(1.5))
txb(s, Inches(6.7), Inches(2.1), Inches(5.9), Inches(0.45),
    "ヒストグラム（棒グラフで度数を表示）", sz=18, bold=True, color=NAVY)
bar_heights = [3, 8, 15, 14, 8, 2]
bar_max = 15
bar_area_h = Inches(3.5)
bar_w = Inches(0.75)
bar_labels = ["40-49", "50-59", "60-69", "70-79", "80-89", "90-99"]
colors_bar = [BLUE, GREEN, ORANGE, ORANGE, GREEN, BLUE]
for i, (h, lbl, c) in enumerate(zip(bar_heights, bar_labels, colors_bar)):
    bh = bar_area_h * h / bar_max
    x = Inches(6.7) + bar_w * i
    by = Inches(2.6) + bar_area_h - bh
    rect(s, x, by, bar_w - Inches(0.05), bh, fill=c, line=WHITE, lw=Pt(0.5))
    txb(s, x, Inches(6.15), bar_w, Inches(0.35),
        lbl, sz=11, color=DKGRAY, align=PP_ALIGN.CENTER)
    txb(s, x, by - Inches(0.3), bar_w, Inches(0.3),
        str(h), sz=14, bold=True, color=DKGRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.38), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35),
    "Excel：ヒストグラム挿入 → [挿入]→[グラフ]→[統計グラフ]→[ヒストグラム]　　Python：plt.hist(scores, bins=6)",
    sz=16, bold=True, color=GOLD)

# ════ S15: Excelで統計計算 ════
s = slide()
header(s, "Excelで統計計算を実践する", "Statistics with Excel")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=GREEN)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.5),
    "Excelの統計関数を使えば、数式を書かずにボタン1つで基本統計量が計算できる",
    sz=20, bold=True, color=WHITE)

txb(s, Inches(0.4), Inches(1.85), Inches(12.5), Inches(0.45),
    "■ よく使うExcel統計関数一覧", sz=19, bold=True, color=NAVY)
th(s, Inches(2.4), [
    ("統計量", Inches(2.5)), ("Excel関数", Inches(3.5)), ("例（A1:A50）", Inches(6.2))
], rh=Inches(0.48))
excel_rows = [
    ("件数 N", "=COUNT(範囲)", "=COUNT(A1:A50) → 50"),
    ("合計 Σxi", "=SUM(範囲)", "=SUM(A1:A50) → 2986"),
    ("算術平均", "=AVERAGE(範囲)", "=AVERAGE(A1:A50) → 59.72"),
    ("中央値", "=MEDIAN(範囲)", "=MEDIAN(A1:A50) → 60"),
    ("最頻値", "=MODE(範囲)", "=MODE(A1:A50) → 62"),
    ("最大値", "=MAX(範囲)", "=MAX(A1:A50) → 83"),
    ("最小値", "=MIN(範囲)", "=MIN(A1:A50) → 41"),
    ("範囲 R", "=MAX()-MIN()", "=MAX(A1:A50)-MIN(A1:A50) → 42"),
]
for i, (stat, func, ex) in enumerate(excel_rows):
    tr(s, Inches(2.88)+Inches(0.51)*i,
       [(stat, Inches(2.5)), (func, Inches(3.5)), (ex, Inches(6.2))],
       rh=Inches(0.49), even=(i%2==0))

rect(s, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35), fill=NAVY)
txb(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
    "「分析ツール」アドイン：[データ]→[データ分析]→[基本統計量] で全統計量を一括出力できる",
    sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════ S16: 次回予告 ════
s = slide()
header(s, "まとめ・次回予告", "Summary & Preview")
rect(s, Inches(0.4), Inches(1.05), Inches(6.0), Inches(5.7), fill=LGRAY, line=BLUE, lw=Pt(1.5))
txb(s, Inches(0.6), Inches(1.15), Inches(5.6), Inches(0.5),
    "■ 第1回で学んだこと", sz=20, bold=True, color=NAVY)
summary = [
    ("Σ記号", "合計を表す数学記号。Σxi = x1+x2+…+xn"),
    ("件数 N", "データの総数（サンプルサイズ）"),
    ("合計", "全データを足し合わせた値"),
    ("範囲 R", "最大値 − 最小値"),
    ("算術平均", "合計 ÷ 件数（外れ値に敏感）"),
    ("幾何平均", "n乗根。成長率・比率の平均に使用"),
    ("中央値", "並べた中央の値（外れ値に強い）"),
    ("最頻値", "最も多く出現する値"),
    ("度数分布表", "階級別に件数を集計した表"),
]
for i, (kw, desc) in enumerate(summary):
    y = Inches(1.75) + Inches(0.55)*i
    rect(s, Inches(0.5), y, Inches(1.8), Inches(0.5), fill=BLUE)
    txb(s, Inches(0.55), y+Inches(0.08), Inches(1.7), Inches(0.42),
        kw, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, Inches(2.35), y+Inches(0.08), Inches(3.85), Inches(0.42),
        desc, sz=15, color=DKGRAY)

rect(s, Inches(6.7), Inches(1.05), Inches(6.2), Inches(5.7), fill=LCYAN, line=BLUE, lw=Pt(1.5))
txb(s, Inches(6.9), Inches(1.15), Inches(5.8), Inches(0.5),
    "■ 次回（第2回）：ばらつきの統計量", sz=20, bold=True, color=NAVY)
next_content = [
    ("分散", "各データと平均の差の2乗の平均", BLUE),
    ("標準偏差", "分散の平方根。ばらつきの基本単位", GREEN),
    ("変動係数", "標準偏差÷平均。異なる単位の比較", ORANGE),
    ("偏差値", "平均50・標準偏差10に変換した指標", PURPLE),
    ("正規分布", "ベル型の確率分布の基礎を学ぶ", RED),
]
for i, (kw, desc, c) in enumerate(next_content):
    y = Inches(1.75) + Inches(0.9)*i
    rect(s, Inches(6.8), y, Inches(1.9), Inches(0.8), fill=c)
    txb(s, Inches(6.9), y+Inches(0.13), Inches(1.7), Inches(0.6),
        kw, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(8.7), y, Inches(4.0), Inches(0.8), fill=WHITE, line=c, lw=Pt(0.5))
    txb(s, Inches(8.8), y+Inches(0.13), Inches(3.8), Inches(0.6),
        desc, sz=16, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.45), fill=NAVY)
txb(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
    "次回までに：Excelでテストデータの平均・中央値・最頻値を実際に計算してみよう",
    sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

out = "C:/Users/saito/Downloads/経営分析/統計学1回_v2.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
