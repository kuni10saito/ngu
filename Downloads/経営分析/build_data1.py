"""
データ経営第1回A_v2.pptx 生成スクリプト
標準本文24pt、全スライドをテキスト/図形で再構築 + シラバス2.0版 2スライド追加
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY    = RGBColor(0x1F, 0x38, 0x64)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x44, 0x72, 0xC4)
GREEN   = RGBColor(0x70, 0xAD, 0x47)
ORANGE  = RGBColor(0xFF, 0xC0, 0x00)
RED     = RGBColor(0xC0, 0x00, 0x00)
PURPLE  = RGBColor(0x70, 0x30, 0xA0)
DKGRAY  = RGBColor(0x26, 0x26, 0x26)
MGRAY   = RGBColor(0x59, 0x59, 0x59)
LGRAY   = RGBColor(0xF2, 0xF2, 0xF2)
LBLUE   = RGBColor(0xBD, 0xD7, 0xEE)
LCYAN   = RGBColor(0xE1, 0xF5, 0xFE)
LORANGE = RGBColor(0xFF, 0xF2, 0xCC)
LGREEN  = RGBColor(0xE2, 0xEF, 0xDA)
LRED    = RGBColor(0xFF, 0xD0, 0xD0)
TEAL    = RGBColor(0x00, 0x70, 0x6E)
LTEAL   = RGBColor(0xD0, 0xF0, 0xEE)
GOLD    = RGBColor(0xB8, 0x86, 0x00)

W = Inches(13.333); H = Inches(7.5)
FONT = "メイリオ"; STD = 24

prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(sl, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    shp = sl.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line: shp.line.color.rgb = line; shp.line.width = lw
    else: shp.line.fill.background()
    return shp

def txb(sl, l, t, w, h, lines, sz=STD, bold=False, color=DKGRAY,
        align=PP_ALIGN.LEFT, wrap=True):
    box = sl.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    if isinstance(lines, str): lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(2)
        run = p.add_run()
        run.text = line; run.font.size = Pt(sz)
        run.font.bold = bold; run.font.color.rgb = color
        run.font.name = FONT
    return box

def header(sl, title, sub=None):
    rect(sl, 0, 0, W, Inches(0.85), fill=NAVY)
    txb(sl, Inches(0.35), Inches(0.08), W-Inches(0.7), Inches(0.72),
        title, sz=26, bold=True, color=WHITE)
    if sub:
        txb(sl, Inches(0.35), Inches(0.92), W-Inches(0.7), Inches(0.5),
            sub, sz=18, bold=True, color=BLUE)

def th(sl, y, cols, rh=Inches(0.55)):
    x = Inches(0.4)
    for txt, w in cols:
        rect(sl, x, y, w, rh, fill=NAVY, line=WHITE, lw=Pt(0.5))
        txb(sl, x+Inches(0.05), y, w-Inches(0.1), rh,
            [txt], sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w

def tr(sl, y, cells, rh=Inches(0.65), even=True):
    bg = LGRAY if even else WHITE; x = Inches(0.4)
    for txt, w, *rest in cells:
        clr = rest[0] if rest else bg
        rect(sl, x, y, w, rh, fill=clr, line=BLUE, lw=Pt(0.5))
        txb(sl, x+Inches(0.05), y, w-Inches(0.1), rh,
            [txt], sz=18, color=DKGRAY, align=PP_ALIGN.CENTER)
        x += w

# ════ S1: タイトル ════
s = slide()
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.05), fill=ORANGE)
txb(s, Inches(1.5), Inches(0.8), Inches(10.3), Inches(1.2),
    "データ経営入門", sz=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.0),
    "第 1 回：データ経営・データサイエンス入門", sz=30, color=ORANGE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.6),
    "齋藤 邦彦", sz=22, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.5),
    "データで経営を変える — データサイエンス×経営学の融合", sz=18, color=LBLUE, align=PP_ALIGN.CENTER)

# ════ S2: 本日の目標 ════
s = slide()
header(s, "本日の目標")
txb(s, Inches(0.6), Inches(1.0), Inches(12.3), Inches(0.5),
    "このセッションを終了する時点で、以下のことができるようになります：",
    sz=20, color=MGRAY)
items = [
    "・ データ経営・データサイエンスの定義と位置付けを説明できる",
    "・ ビッグデータ時代の背景と重要性を理解できる",
    "・ データドリブン経営の具体的な効果・事例を挙げられる",
    "・ DX時代のデータ利活用の2つの柱（分析支援・業務自動化）を説明できる",
    "・ 産業別（金融・工場・農漁業・商店）のデータ活用を理解できる",
    "・ 生成AI・LLMのビジネス活用を概説できる",
    "・ 【New】Claude Code でデータ経営分析を自動化できる",
]
for i, item in enumerate(items):
    clr = LCYAN if i == len(items)-1 else LGRAY
    rect(s, Inches(0.5), Inches(1.6)+Inches(0.7)*i, Inches(12.3), Inches(0.62), fill=clr, line=BLUE, lw=Pt(0.5))
    c = TEAL if i == len(items)-1 else DKGRAY
    b = i == len(items)-1
    txb(s, Inches(0.65), Inches(1.65)+Inches(0.7)*i, Inches(12.0), Inches(0.6),
        item, sz=20, bold=b, color=c)

# ════ S3: データ経営とは ════
s = slide()
header(s, "データ経営とは", "Data-Driven Management")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(1.2), fill=NAVY)
txb(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.1),
    "データを用いた経営・データによる経営学", sz=28, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.6),
    "データサイエンスを基礎とし、経営的課題を解決するための意思決定手法・プロセス", sz=20, color=LBLUE)

cols_data = [
    ("概念", "データ経営は、データサイエンスを用いてビジネスの意思決定を科学的に支援する経営手法", LCYAN),
    ("対象", "マーケティング、ファイナンス、会計、人事など経営の広い応用分野", LGREEN),
    ("目的", "意思決定の質・速度を高め、競合優位性を確立する", LORANGE),
    ("手段", "統計分析・機械学習・AI・可視化ツールを組み合わせてデータを価値に変換", LBLUE),
]
y = Inches(2.4)
for i, (lbl, desc, bg) in enumerate(cols_data):
    rect(s, Inches(0.4), y, Inches(2.8), Inches(0.85), fill=NAVY)
    txb(s, Inches(0.5), y+Inches(0.1), Inches(2.6), Inches(0.75),
        lbl, sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.2), y, Inches(9.7), Inches(0.85), fill=bg)
    txb(s, Inches(3.3), y+Inches(0.08), Inches(9.5), Inches(0.8),
        desc, sz=20, color=DKGRAY)
    y += Inches(0.95)

txb(s, Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.5),
    "📌 データ経営 ＝ 統計・AI・プログラミングの力で経営判断を科学する学問",
    sz=20, bold=True, color=NAVY)

# ════ S4: データサイエンスとは ════
s = slide()
header(s, "データサイエンスとは", "Data Science Foundation")
# 3つの柱
titles = ["統計学・数学", "情報技術", "ビジネス・経営"]
colors = [BLUE, GREEN, ORANGE]
descs = [
    ["・ 記述統計・推測統計", "・ 確率論・線形代数", "・ 機械学習理論"],
    ["・ Python・Excel", "・ データベース・SQL", "・ クラウド・AI API"],
    ["・ 業務課題の定義", "・ KPI設計・意思決定", "・ 経営戦略への応用"],
]
for i in range(3):
    x = Inches(0.5) + Inches(4.2)*i
    rect(s, x, Inches(1.1), Inches(4.0), Inches(0.7), fill=colors[i])
    txb(s, x+Inches(0.1), Inches(1.15), Inches(3.8), Inches(0.65),
        titles[i], sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(1.8), Inches(4.0), Inches(2.2), fill=LGRAY, line=colors[i], lw=Pt(1.5))
    txb(s, x+Inches(0.15), Inches(1.9), Inches(3.7), Inches(2.0),
        descs[i], sz=20, color=DKGRAY)

# 重なり部分
rect(s, Inches(4.2), Inches(4.2), Inches(4.9), Inches(1.0), fill=LCYAN, line=NAVY, lw=Pt(2))
txb(s, Inches(4.3), Inches(4.3), Inches(4.7), Inches(0.9),
    "データサイエンティスト\nの活躍領域", sz=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txb(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
    "データサイエンスの特徴：", sz=20, bold=True, color=NAVY)
txb(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9),
    ["・ 課題解決に向けたデータ処理・分析の手法体系（学際的）",
     "・ 様々な分野（科学・金融・ヘルスケア）に適用可能  ・ AIと組み合わせで自動化"],
    sz=20, color=DKGRAY)

# ════ S5: データ経営で学ぶべき内容 ════
s = slide()
header(s, "データ経営：学ぶべき内容", "Learning Roadmap")
categories = [
    ("データサイエンス基礎", BLUE,
     ["データ収集・クリーニング", "データ処理・分析手法", "データ可視化（グラフ・ダッシュボード）"]),
    ("統計学・数学", GREEN,
     ["記述統計・推測統計", "回帰分析・相関分析", "機械学習・ディープラーニング"]),
    ("プログラミング・ツール", ORANGE,
     ["Excel（関数・ピボット）", "Python（pandas・sklearn）", "生成AI・LLM・ChatGPT活用"]),
    ("経営・ビジネス応用", PURPLE,
     ["経営意思決定支援", "マーケティング分析", "ファイナンス・会計データ分析"]),
]
for i, (cat, color, items) in enumerate(categories):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(1.1), Inches(3.0), Inches(0.65), fill=color)
    txb(s, x+Inches(0.1), Inches(1.15), Inches(2.8), Inches(0.6),
        cat, sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(1.75), Inches(3.0), Inches(2.5), fill=LGRAY, line=color, lw=Pt(1))
    for j, item in enumerate(items):
        txb(s, x+Inches(0.15), Inches(1.9)+Inches(0.75)*j, Inches(2.75), Inches(0.65),
            f"▶ {item}", sz=19, color=DKGRAY)

rect(s, Inches(0.4), Inches(4.45), Inches(12.5), Inches(0.65), fill=NAVY)
txb(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
    "本講義のゴール：データで経営判断を科学できる「データ経営人材」の育成",
    sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(s, Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.5),
    "この講義で使うツール・技術：", sz=20, bold=True, color=NAVY)
tools = [
    ("Excel", LBLUE, "基本分析・可視化"),
    ("Python", LGREEN, "機械学習・自動化"),
    ("ChatGPT/Claude", LORANGE, "生成AI活用"),
    ("EDINET API", LTEAL, "財務データ取得"),
]
for i, (tool, bg, desc) in enumerate(tools):
    x = Inches(0.4) + Inches(3.1)*i
    rect(s, x, Inches(5.9), Inches(3.0), Inches(1.3), fill=bg, line=BLUE, lw=Pt(0.8))
    txb(s, x+Inches(0.1), Inches(6.0), Inches(2.8), Inches(0.5),
        tool, sz=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txb(s, x+Inches(0.1), Inches(6.5), Inches(2.8), Inches(0.55),
        desc, sz=17, color=DKGRAY, align=PP_ALIGN.CENTER)

# ════ S6: ビッグデータとデータサイエンス ════
s = slide()
header(s, "ビッグデータの時代とデータサイエンス", "Big Data Era")
# 3V定義
vs = [("Volume\n大量性", "SNS・IoT・購買履歴など\n膨大なデータが蓄積", BLUE),
      ("Variety\n多様性", "テキスト・画像・音声・\nセンサーなど多種多様", GREEN),
      ("Velocity\n即時性", "リアルタイムで生成・\n更新されるデータ", ORANGE)]
for i, (v, desc, c) in enumerate(vs):
    x = Inches(0.5) + Inches(2.8)*i
    rect(s, x, Inches(1.1), Inches(2.6), Inches(1.5), fill=c)
    txb(s, x+Inches(0.1), Inches(1.2), Inches(2.4), Inches(1.3),
        v, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.6), Inches(2.6), Inches(1.1), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(2.7), Inches(2.4), Inches(1.0),
        desc, sz=17, color=DKGRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(9.0), Inches(1.1), Inches(3.9), Inches(2.6), fill=LCYAN, line=NAVY, lw=Pt(2))
txb(s, Inches(9.1), Inches(1.2), Inches(3.7), Inches(0.6),
    "ビッグデータの源泉", sz=20, bold=True, color=NAVY)
sources = ["・ 人々の活動（購買履歴・SNS・Web閲覧）",
           "・ 科学分野（人工衛星・GPS観測）",
           "・ IoTセンサー（工場・車・家電）",
           "・ 金融取引データ（高頻度取引）"]
txb(s, Inches(9.1), Inches(1.85), Inches(3.7), Inches(1.8),
    sources, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.6), fill=NAVY)
txb(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.55),
    "データサイエンス：社会に遍在するあらゆる種類のデータを処理・分析して有用な情報（価値）を引き出す技術",
    sz=19, bold=True, color=WHITE)

rows = [
    ("1980-90年代", "データ蓄積開始", "データベース・ERPの普及"),
    ("2000年代", "ビッグデータ台頭", "SNS・EC・IoTの爆発的拡大"),
    ("2010年代", "機械学習実用化", "ディープラーニング・クラウド活用"),
    ("2020年代", "生成AI時代", "LLM・マルチモーダルAIの普及"),
]
th(s, Inches(4.75), [("時代", Inches(2.5)), ("段階", Inches(3.5)), ("特徴", Inches(6.2))], rh=Inches(0.5))
for i, (a, b, c) in enumerate(rows):
    tr(s, Inches(5.25)+Inches(0.55)*i, [(a, Inches(2.5)), (b, Inches(3.5)), (c, Inches(6.2))],
       rh=Inches(0.52), even=(i%2==0))

# ════ S7: 資源としてのデータ ════
s = slide()
header(s, "資源としてのデータ", "Data as the New Oil")
rect(s, Inches(0.4), Inches(1.05), Inches(5.8), Inches(5.2), fill=LORANGE, line=ORANGE, lw=Pt(2))
txb(s, Inches(0.6), Inches(1.15), Inches(5.4), Inches(0.7),
    "データ ＝ 21世紀の石油", sz=24, bold=True, color=GOLD)
txb(s, Inches(0.6), Inches(1.9), Inches(5.4), Inches(1.3),
    ["石油：採掘して初めて価値が生まれる",
     "データ：収集・分析して初めて価値が生まれる",
     "→ 適切な「加工技術」が競争優位を生む"],
    sz=19, color=DKGRAY)
txb(s, Inches(0.6), Inches(3.3), Inches(5.4), Inches(0.5),
    "データサイエンティストの役割", sz=20, bold=True, color=NAVY)
roles = [
    ("情報学", "データをコンピュータで処理する技術"),
    ("統計学", "データの分析に関する知識・手法"),
    ("経営学", "応用分野の知識（人文科学・経営）"),
]
for i, (role, desc) in enumerate(roles):
    rect(s, Inches(0.6), Inches(3.9)+Inches(0.65)*i, Inches(1.8), Inches(0.58), fill=NAVY)
    txb(s, Inches(0.65), Inches(3.95)+Inches(0.65)*i, Inches(1.7), Inches(0.52),
        role, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.4), Inches(3.9)+Inches(0.65)*i, Inches(3.6), Inches(0.58), fill=WHITE, line=ORANGE, lw=Pt(0.5))
    txb(s, Inches(2.5), Inches(3.95)+Inches(0.65)*i, Inches(3.4), Inches(0.52),
        desc, sz=18, color=DKGRAY)

txb(s, Inches(0.5), Inches(5.95), Inches(5.8), Inches(0.4),
    "→ 文理融合的な人材が求められる", sz=19, bold=True, color=RED)

rect(s, Inches(6.5), Inches(1.05), Inches(6.4), Inches(5.2), fill=LGRAY, line=BLUE, lw=Pt(1.5))
txb(s, Inches(6.7), Inches(1.15), Inches(6.0), Inches(0.6),
    "データの価値化プロセス", sz=22, bold=True, color=NAVY)
steps = [
    ("① 収集", "IoT・API・スクレイピング・調査", BLUE),
    ("② 蓄積", "データベース・データレイク・クラウド", GREEN),
    ("③ 処理", "クリーニング・前処理・特徴量エンジニアリング", ORANGE),
    ("④ 分析", "統計分析・機械学習・AI", PURPLE),
    ("⑤ 活用", "意思決定・可視化・自動化・新サービス", RED),
]
for i, (step, desc, c) in enumerate(steps):
    rect(s, Inches(6.6), Inches(1.85)+Inches(0.82)*i, Inches(1.3), Inches(0.72), fill=c)
    txb(s, Inches(6.65), Inches(1.9)+Inches(0.82)*i, Inches(1.2), Inches(0.65),
        step, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(7.9), Inches(1.85)+Inches(0.82)*i, Inches(4.8), Inches(0.72), fill=WHITE, line=c, lw=Pt(0.5))
    txb(s, Inches(8.0), Inches(1.9)+Inches(0.82)*i, Inches(4.6), Inches(0.65),
        desc, sz=18, color=DKGRAY)

# ════ S8: DX時代のデータ利活用 ════
s = slide()
header(s, "DX時代のデータ利活用", "Data Utilization in DX Era")
txb(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.45),
    "データ活用の2大柱：意思決定支援 × 業務自動化", sz=21, bold=True, color=NAVY)

rect(s, Inches(0.4), Inches(1.55), Inches(6.2), Inches(4.9), fill=LCYAN, line=BLUE, lw=Pt(2))
txb(s, Inches(0.5), Inches(1.65), Inches(6.0), Inches(0.65),
    "(1) 意思決定支援（レポーティング中心）", sz=21, bold=True, color=NAVY)
items1 = [
    "■ 主な手法",
    "・ ExcelやBIツールを使った統計分析",
    "・ 現状把握・原因仮説の探索・効果測定",
    "・ 「PDCA」的な流れに沿った定量的判断支援",
    "",
    "■ 担当者",
    "・ データアナリストが中心",
    "・ 業務部門と連携した分析推進",
    "",
    "■ 出力物",
    "・ ダッシュボード・月次レポート・KPI管理表",
]
txb(s, Inches(0.55), Inches(2.35), Inches(5.9), Inches(3.8),
    items1, sz=19, color=DKGRAY)

rect(s, Inches(6.9), Inches(1.55), Inches(6.0), Inches(4.9), fill=LGREEN, line=GREEN, lw=Pt(2))
txb(s, Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.65),
    "(2) モデル構築・業務自動化（生産性向上）", sz=21, bold=True, color=TEAL)
items2 = [
    "■ 主な手法",
    "・ AI・機械学習でモデル構築",
    "・ 業務システムへの組み込み・自動化",
    "・ マーケティングや需要予測に活用",
    "",
    "■ 担当者",
    "・ データサイエンティストが担当",
    "・ エンジニアと連携した実装推進",
    "",
    "■ 出力物",
    "・ 予測モデル・推薦システム・異常検知",
]
txb(s, Inches(7.05), Inches(2.35), Inches(5.7), Inches(3.8),
    items2, sz=19, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.65), fill=NAVY)
txb(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.58),
    "両者を兼ね備えた人材が「データ経営」を実現する ← 本講義で目指すゴール",
    sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════ S9: データドリブン意思決定の重要性 ════
s = slide()
header(s, "データドリブン意思決定の重要性", "Evidence-Based Management")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(1.0), fill=LCYAN, line=BLUE, lw=Pt(1))
txb(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5),
    "✅ データリテラシーと業績指標は「正の相関」がある", sz=22, bold=True, color=NAVY)
txb(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.45),
    "「データは21世紀の石油」— データを分析してビジネスに活かす能力が企業の競争力を左右する",
    sz=18, color=DKGRAY)

evidences = [
    ("Erik Brynjolfsson\n(2011年・MIT)",
     "データドリブンな企業は他社に比べて生産性が5〜6%高い\n資産活用・ROI・市場価値においても正の相関あり",
     BLUE),
    ("クリッテック・ジャパン\n(2018年)",
     "「データリテラシー指数」を発表\n社員のデータ活用スキルと業績の正の相関を定量化",
     GREEN),
    ("McKinsey Global\nInstitute",
     "データドリブン企業の収益成長率は業界平均の23倍\n顧客獲得コストが6分の1以下に低減",
     ORANGE),
]
for i, (src, desc, c) in enumerate(evidences):
    x = Inches(0.4) + Inches(4.2)*i
    rect(s, x, Inches(2.2), Inches(4.0), Inches(0.7), fill=c)
    txb(s, x+Inches(0.1), Inches(2.25), Inches(3.8), Inches(0.65),
        src, sz=17, bold=True, color=WHITE)
    rect(s, x, Inches(2.9), Inches(4.0), Inches(1.5), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(2.95), Inches(3.8), Inches(1.4),
        desc, sz=17, color=DKGRAY)

txb(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.45),
    "✅ データ分析スキルはビジネス成果に不可欠", sz=20, bold=True, color=RED)
txb(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.5),
    "組織内でデータ利活用スキルが不足していると、ビジネス成果を年々増やすことが困難",
    sz=19, color=DKGRAY)

rows2 = [
    ("データ活用レベル", "意思決定の質", "競争優位性", "代表企業例"),
    ("レベル1：経験則中心", "主観的・属人的", "低", "多くの中小企業"),
    ("レベル2：Excel分析", "定量的・周期的", "中", "一般的な中堅企業"),
    ("レベル3：BI・ML活用", "予測的・自動的", "高", "DX先進企業"),
    ("レベル4：AI駆動経営", "リアルタイム最適化", "最高", "GAFAM・メガテック"),
]
th(s, Inches(5.75), [("データ活用レベル", Inches(3.0)), ("意思決定の質", Inches(3.0)),
                      ("競争優位性", Inches(2.5)), ("代表企業例", Inches(3.7))], rh=Inches(0.48))
bg_list = [LGRAY, LCYAN, LGREEN, LORANGE]
for i, row in enumerate(rows2[1:]):
    tr(s, Inches(6.23)+Inches(0.52)*i,
       [(row[0], Inches(3.0)), (row[1], Inches(3.0)), (row[2], Inches(2.5)), (row[3], Inches(3.7))],
       rh=Inches(0.5), even=(i%2==0))

# ════ S10: データサイエンスの産業別利用例 ════
s = slide()
header(s, "データサイエンスの産業別利用例", "Industry Applications")
industries = [
    ("ヘルスケア", BLUE, [
        "・ 病気の早期発見・診断支援（画像AI）",
        "・ 新薬開発の候補化合物予測",
        "・ 電子カルテのデータ分析",
        "・ 患者リスク層別化モデル",
    ]),
    ("金融", GREEN, [
        "・ 信用リスク評価・不正行為検出",
        "・ アルゴリズム取引・高頻度取引",
        "・ 顧客セグメンテーション",
        "・ 投資戦略の最適化",
    ]),
    ("製造・工場", ORANGE, [
        "・ スマートファクトリー（IoT×AI）",
        "・ 故障予測・予防保全",
        "・ 工程の総合管理・最適化",
        "・ 品質検査の自動化",
    ]),
    ("農漁業", PURPLE, [
        "・ ハウス農業（温度・CO2管理）",
        "・ 漁場予測（海水温画像処理）",
        "・ 需要予測による廃棄ロス削減",
        "・ ドローン・IoTによる圃場管理",
    ]),
    ("商店・流通", RED, [
        "・ 売れ筋分析（基本統計・可視化）",
        "・ 需要予測・在庫最適化",
        "・ 顧客行動分析・推薦システム",
        "・ 価格最適化（ダイナミックプライシング）",
    ]),
    ("交通・都市", TEAL, [
        "・ 交通流予測・渋滞対策",
        "・ 公共交通スケジュール最適化",
        "・ 自動運転（センサーデータ統合）",
        "・ スマートシティ・エネルギー管理",
    ]),
]
for i, (ind, c, items) in enumerate(industries):
    col = i % 3; row = i // 3
    x = Inches(0.4) + Inches(4.3)*col
    y = Inches(1.1) + Inches(2.85)*row
    rect(s, x, y, Inches(4.1), Inches(0.6), fill=c)
    txb(s, x+Inches(0.1), y+Inches(0.08), Inches(3.9), Inches(0.55),
        ind, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, y+Inches(0.6), Inches(4.1), Inches(2.1), fill=LGRAY, line=c, lw=Pt(0.8))
    txb(s, x+Inches(0.1), y+Inches(0.7), Inches(3.9), Inches(2.0),
        items, sz=17, color=DKGRAY)

# ════ S11: データサイエンス×経営学 ════
s = slide()
header(s, "データサイエンス × 経営学 ＝ データ経営", "Data Science × Business")
rect(s, Inches(0.4), Inches(1.05), Inches(5.5), Inches(5.6), fill=LCYAN, line=BLUE, lw=Pt(2))
txb(s, Inches(0.5), Inches(1.15), Inches(5.3), Inches(0.65),
    "データサイエンス", sz=24, bold=True, color=BLUE)
ds_items = ["・ データ収集・前処理・分析", "・ 統計学・機械学習", "・ 可視化・モデリング",
            "・ Python・SQL・BI活用", "・ ビッグデータ処理技術"]
txb(s, Inches(0.6), Inches(1.9), Inches(5.1), Inches(2.5),
    ds_items, sz=20, color=DKGRAY)
rect(s, Inches(0.5), Inches(4.5), Inches(5.2), Inches(1.8), fill=BLUE)
txb(s, Inches(0.6), Inches(4.6), Inches(5.0), Inches(1.6),
    ["適用領域：科学・金融・医療・製造",
     "目的：データから価値・知識を引き出す",
     "出力：予測モデル・分析レポート"],
    sz=18, color=WHITE)

rect(s, Inches(6.8), Inches(1.05), Inches(6.1), Inches(5.6), fill=LGREEN, line=GREEN, lw=Pt(2))
txb(s, Inches(6.9), Inches(1.15), Inches(5.9), Inches(0.65),
    "経営学", sz=24, bold=True, color=TEAL)
biz_items = ["・ 経営戦略・マーケティング", "・ ファイナンス・会計分析",
             "・ 組織・人事・オペレーション", "・ 意思決定プロセス", "・ 競争優位の構築"]
txb(s, Inches(7.0), Inches(1.9), Inches(5.7), Inches(2.5),
    biz_items, sz=20, color=DKGRAY)
rect(s, Inches(6.9), Inches(4.5), Inches(5.8), Inches(1.8), fill=TEAL)
txb(s, Inches(7.0), Inches(4.6), Inches(5.6), Inches(1.6),
    ["適用領域：企業・組織・市場",
     "目的：経営判断・競争力向上",
     "出力：戦略提言・KPI管理・利益最大化"],
    sz=18, color=WHITE)

rect(s, Inches(5.7), Inches(2.5), Inches(1.7), Inches(2.0), fill=GOLD)
txb(s, Inches(5.8), Inches(2.6), Inches(1.5), Inches(1.8),
    "データ\n経営\n（融合）", sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.55), fill=NAVY)
txb(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.48),
    "AI経営 ≈ データドリブン経営 + 自動化　— 生成AI時代の経営の新潮流",
    sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════ S12: 生成AIとLLM ════
s = slide()
header(s, "生成AIとLLM（大規模言語モデル）", "Generative AI & LLM")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(1.0), fill=NAVY)
txb(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5),
    "生成AIとは：既存のデータに基づいてテキスト・画像・音声・動画などを生成する技術", sz=20, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.4),
    "LLM（Large Language Model）= テキストを自動生成する大規模言語モデル", sz=20, color=LBLUE)

types = [
    ("テキスト生成", BLUE, ["ChatGPT / Claude / Gemini", "文書作成・要約・翻訳", "コード生成・デバッグ"]),
    ("画像生成", GREEN, ["DALL-E 3 / Midjourney", "Stable Diffusion", "デザイン・マーケティング素材"]),
    ("音声合成", ORANGE, ["ElevenLabs / VOICEVOX", "テキスト→音声変換(TTS)", "音声クローン・ナレーション"]),
    ("動画生成", PURPLE, ["Sora / Runway / Kling", "テキスト→動画生成", "広告・教育コンテンツ"]),
]
for i, (t, c, items) in enumerate(types):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(2.15), Inches(3.05), Inches(0.65), fill=c)
    txb(s, x+Inches(0.1), Inches(2.2), Inches(2.85), Inches(0.6),
        t, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.8), Inches(3.05), Inches(1.8), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(2.9), Inches(2.85), Inches(1.7),
        items, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(4.75), Inches(12.5), Inches(0.55), fill=TEAL)
txb(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5),
    "経営×生成AI：業務自動化・意思決定支援・新サービス開発に広く普及中",
    sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

biz_apps = [
    ("営業・マーケ", ["顧客提案書の自動生成", "SNS投稿の最適化", "需要予測レポート作成"]),
    ("経営企画", ["財務分析レポート自動化", "競合分析・SWOT生成", "会議議事録の要約"]),
    ("人事・採用", ["求人票・評価シート生成", "研修コンテンツ作成", "面接質問の最適化"]),
    ("開発・IT", ["コード生成・レビュー", "テスト自動化", "ドキュメント生成"]),
]
for i, (dept, items) in enumerate(biz_apps):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(5.45), Inches(3.05), Inches(0.55), fill=NAVY)
    txb(s, x+Inches(0.1), Inches(5.5), Inches(2.85), Inches(0.5),
        dept, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(6.0), Inches(3.05), Inches(1.3), fill=LGRAY, line=NAVY, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(6.05), Inches(2.85), Inches(1.2),
        items, sz=16, color=DKGRAY)

# ════ S13: まとめ ════
s = slide()
header(s, "第1回 まとめ", "Lecture Summary")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.5), fill=BLUE)
txb(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.45),
    "今日学んだキーワード", sz=20, bold=True, color=WHITE)

keywords = [
    ("データ経営", "データサイエンスを基礎とした科学的経営手法。意思決定の質を高める"),
    ("データサイエンス", "統計・IT・ビジネスの3分野を統合した学際的問題解決技術"),
    ("ビッグデータ", "3V（Volume・Variety・Velocity）を特徴とする大量多様なデータ"),
    ("DXデータ活用", "意思決定支援（分析）× 業務自動化（AI/ML）の2つの柱"),
    ("生成AI・LLM", "テキスト・画像・音声を生成するAI。経営全般に浸透中"),
    ("産業別応用", "ヘルスケア・金融・工場・農漁業・商店それぞれのデータ活用"),
]
for i, (kw, desc) in enumerate(keywords):
    y = Inches(1.65) + Inches(0.82)*i
    rect(s, Inches(0.4), y, Inches(2.8), Inches(0.72), fill=NAVY)
    txb(s, Inches(0.5), y+Inches(0.1), Inches(2.6), Inches(0.62),
        kw, sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.2), y, Inches(9.7), Inches(0.72), fill=LGRAY, line=BLUE, lw=Pt(0.3))
    txb(s, Inches(3.3), y+Inches(0.1), Inches(9.5), Inches(0.62),
        desc, sz=19, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.65), fill=LORANGE, line=ORANGE, lw=Pt(1))
txb(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.55),
    "次回：データ収集・前処理の実践（Excel / Python）と基本統計量の計算",
    sz=20, bold=True, color=GOLD)

# ════ S14: 課題 ════
s = slide()
header(s, "課題・演習", "Assignment")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=ORANGE)
txb(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6),
    "【課題】日常で使われるデータを考え、簡単なデータ分析を行いプレゼン資料を作成する",
    sz=21, bold=True, color=WHITE)

rect(s, Inches(0.4), Inches(1.85), Inches(12.5), Inches(2.5), fill=LGRAY, line=BLUE, lw=Pt(1))
txb(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
    "■ 身近なデータ例（いずれかを選択）", sz=20, bold=True, color=NAVY)
examples = [
    "・ アルバイト先のハンバーガー店の売上（時間帯別・日別・月別）",
    "・ 通学に使う駅の乗降者数（平日・週間・季節変動）",
    "・ 自分のお小遣い・出費の記録",
    "・ 好きなスポーツチームの成績データ",
    "・ SNSのフォロワー数・いいね数の推移",
]
txb(s, Inches(0.6), Inches(2.5), Inches(12.0), Inches(1.7),
    examples, sz=19, color=DKGRAY)

rect(s, Inches(0.4), Inches(4.5), Inches(6.0), Inches(2.7), fill=LCYAN, line=BLUE, lw=Pt(1))
txb(s, Inches(0.5), Inches(4.6), Inches(5.8), Inches(0.55),
    "■ プレゼン資料の構成", sz=20, bold=True, color=NAVY)
steps2 = ["① 選んだデータの説明（何のデータか）",
          "② データの収集方法",
          "③ データの可視化（グラフ・表）",
          "④ わかったこと・気づき",
          "⑤ 今後どう活用できるか"]
txb(s, Inches(0.6), Inches(5.2), Inches(5.7), Inches(1.8),
    steps2, sz=18, color=DKGRAY)

rect(s, Inches(6.8), Inches(4.5), Inches(6.1), Inches(2.7), fill=LORANGE, line=ORANGE, lw=Pt(1))
txb(s, Inches(6.9), Inches(4.6), Inches(5.9), Inches(0.55),
    "■ 提出方法・期限", sz=20, bold=True, color=GOLD)
txb(s, Inches(6.9), Inches(5.2), Inches(5.9), Inches(1.8),
    ["・ 今日は提出なし（次回から）",
     "・ PowerPoint または PDF で提出",
     "・ 5〜8スライドでまとめる",
     "・ データソースを明記すること"],
    sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(7.25), Inches(12.5), Inches(0.05), fill=ORANGE)

# ════ S15: 国内データ経営成功事例 ════
s = slide()
header(s, "国内企業のデータ経営成功事例", "Japanese Data-Driven Management Cases")
cases_jp = [
    ("トヨタ自動車", BLUE, "製造×IoT×データ",
     ["工場の全設備をセンサーで監視",
      "品質不良を0.01%未満に削減",
      "需要予測AIで在庫を30%削減",
      "TPS×データで世界最高水準の生産性"]),
    ("楽天グループ", GREEN, "EC×データ活用",
     ["購買履歴×行動ログで超個人化",
      "レコメンドAIがCVR3倍向上",
      "楽天エコシステムで横断分析",
      "データ共有で顧客LTV最大化"]),
    ("NTTドコモ", ORANGE, "通信×ビッグデータ",
     ["位置情報データで人流を可視化",
      "コロナ禍の行動変容を行政に提供",
      "モバイル空間統計で都市計画支援",
      "通信データを社会インフラとして活用"]),
    ("ファーストリテイリング", PURPLE, "小売×需要予測",
     ["RFID全商品タグで在庫リアルタイム管理",
      "需要予測AIで廃棄ロスを大幅削減",
      "AIデザインアシスタントで企画効率化",
      "サプライチェーン全体をデジタル化"]),
]
for i, (co, c, tag, items) in enumerate(cases_jp):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(1.05), Inches(3.05), Inches(0.65), fill=c)
    txb(s, x+Inches(0.08), Inches(1.1), Inches(2.9), Inches(0.35),
        co, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, x+Inches(0.08), Inches(1.45), Inches(2.9), Inches(0.25),
        tag, sz=13, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(1.7), Inches(3.05), Inches(3.0), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(1.8), Inches(2.85), Inches(2.8),
        items, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.5), fill=LCYAN, line=BLUE, lw=Pt(1))
txb(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.4),
    "共通点：①データ収集基盤への投資 ②全社的なデータ活用文化 ③AI/機械学習の継続的活用 ④経営判断へのフィードバック",
    sz=17, bold=True, color=NAVY)
txb(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.45),
    "■ 中小企業でもできるデータ経営の第一歩", sz=19, bold=True, color=NAVY)
smb_steps = [
    ("STEP 1", "既存データの棚卸し（売上・顧客・在庫）", BLUE),
    ("STEP 2", "Excelで基本集計とグラフ化を習慣化", GREEN),
    ("STEP 3", "KPIを定義して定期的にモニタリング", ORANGE),
    ("STEP 4", "データに基づく仮説→実験→改善サイクル", PURPLE),
]
for i, (step, desc, c) in enumerate(smb_steps):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(6.05), Inches(3.05), Inches(0.45), fill=c)
    txb(s, x+Inches(0.08), Inches(6.1), Inches(2.9), Inches(0.4),
        step, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(6.5), Inches(3.05), Inches(0.7), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.08), Inches(6.55), Inches(2.9), Inches(0.6),
        desc, sz=15, color=DKGRAY, align=PP_ALIGN.CENTER)

# ════ S16: DXと経営変革 ════
s = slide()
header(s, "DX（デジタルトランスフォーメーション）と経営変革", "Digital Transformation")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.75), fill=NAVY)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.4),
    "DX = デジタル技術を活用してビジネスモデル・組織・文化を根本から変革すること",
    sz=21, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.53), Inches(12.1), Inches(0.25),
    "単なる「IT化」や「システム導入」ではなく、価値創造の方法そのものを変えること", sz=17, color=LBLUE)

dx_levels = [
    ("デジタイゼーション", MGRAY, "アナログ→デジタル変換",
     "紙の書類をExcelに入力・FAXをメールに変更\n「作業の効率化」レベル"),
    ("デジタライゼーション", BLUE, "業務プロセスの変革",
     "販売管理をクラウドSaaSで自動化・顧客DBとの連携\n「業務の再設計」レベル"),
    ("デジタルトランスフォーメーション", PURPLE, "ビジネスモデルの変革",
     "データを核にした新サービス・収益モデルの創出\n「企業価値の再定義」レベル"),
]
for i, (name, c, sub, desc) in enumerate(dx_levels):
    y = Inches(2.0) + Inches(1.4)*i
    rect(s, Inches(0.4), y, Inches(12.5), Inches(1.25), fill=LGRAY if i%2==0 else WHITE, line=c, lw=Pt(1.5))
    rect(s, Inches(0.4), y, Inches(3.5), Inches(1.25), fill=c)
    txb(s, Inches(0.5), y+Inches(0.1), Inches(3.3), Inches(0.55),
        name, sz=17, bold=True, color=WHITE)
    txb(s, Inches(0.5), y+Inches(0.65), Inches(3.3), Inches(0.5),
        sub, sz=14, color=WHITE)
    txb(s, Inches(4.1), y+Inches(0.15), Inches(8.5), Inches(1.0),
        desc, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.5), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.42),
    "経産省「2025年の崖」：DXに失敗した企業は最大12兆円/年の経済損失と推計。既存IT維持コストがDX投資を圧迫",
    sz=18, bold=True, color=GOLD)
txb(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.42),
    "→ データ経営の構築はDX推進の核心。「データを持つ企業」と「持たない企業」の格差は急拡大している",
    sz=18, color=DKGRAY)

# ════ S17: データ組織の構築 ════
s = slide()
header(s, "データ組織の構築とデータ人材", "Data Organization & Talent")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=TEAL)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.5),
    "データ経営を推進するには「人材」「組織」「文化」の3つを同時に整備する必要がある",
    sz=20, bold=True, color=WHITE)

roles = [
    ("CDO\nChief Data Officer", NAVY, "最高データ責任者",
     ["データ戦略の立案・推進",
      "データガバナンスの統括",
      "経営陣とのデータ活用調整",
      "組織横断のデータ文化醸成"]),
    ("データサイエンティスト", BLUE, "分析・モデル開発",
     ["機械学習モデルの開発",
      "統計分析・仮説検証",
      "ビジネス課題の数値化",
      "分析結果の経営への提言"]),
    ("データエンジニア", GREEN, "基盤・パイプライン",
     ["データ収集・加工・蓄積",
      "データウェアハウス構築",
      "ETLパイプラインの開発",
      "データ品質の維持管理"]),
    ("ビジネスアナリスト", ORANGE, "分析・意思決定支援",
     ["KPI設計・ダッシュボード",
      "現場データの収集・整理",
      "SQLによるデータ抽出",
      "レポート作成・説明"]),
]
for i, (title, c, sub, items) in enumerate(roles):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(1.9), Inches(3.05), Inches(0.8), fill=c)
    txb(s, x+Inches(0.08), Inches(1.95), Inches(2.9), Inches(0.45),
        title, sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, x+Inches(0.08), Inches(2.4), Inches(2.9), Inches(0.25),
        sub, sz=13, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.7), Inches(3.05), Inches(2.75), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(2.8), Inches(2.85), Inches(2.55),
        items, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.5), fill=LCYAN, line=BLUE)
txb(s, Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.4),
    "CoE（Center of Excellence）: 全社データ活用の専門チームを設置し、各部門にデータ活用を浸透させる体制",
    sz=18, bold=True, color=NAVY)
txb(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.45),
    "■ データ人材の不足：経済産業省は2030年に約79万人のIT・データ人材が不足すると試算", sz=18, bold=True, color=RED)
txb(s, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.45),
    "→ この講義でデータ分析の基礎を身につけることが、キャリアの大きな強みになる",
    sz=18, color=DKGRAY)

# ════ S18: データガバナンス ════
s = slide()
header(s, "データガバナンスとセキュリティ", "Data Governance & Security")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=RED)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.5),
    "データガバナンス = データを「正しく・安全に・効果的に」活用するためのルールと体制",
    sz=20, bold=True, color=WHITE)

gov_pillars = [
    ("データ品質", RED,
     ["正確性：データの誤りを排除",
      "完全性：欠損・漏れをなくす",
      "一貫性：複数DBで矛盾しない",
      "最新性：リアルタイムで更新"]),
    ("セキュリティ", NAVY,
     ["アクセス制御（権限管理）",
      "暗号化・マスキング処理",
      "不正アクセス検知・ログ管理",
      "インシデント対応手順の整備"]),
    ("プライバシー保護", PURPLE,
     ["個人情報保護法への対応",
      "GDPR（EU規制）の理解",
      "同意管理・データ削除権",
      "仮名化・匿名化処理"]),
    ("データカタログ", GREEN,
     ["データ資産の一元管理",
      "メタデータの整備・検索",
      "データオーナーの明確化",
      "データ利活用ルールの文書化"]),
]
for i, (title, c, items) in enumerate(gov_pillars):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(1.9), Inches(3.05), Inches(0.6), fill=c)
    txb(s, x+Inches(0.08), Inches(1.95), Inches(2.9), Inches(0.55),
        title, sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.5), Inches(3.05), Inches(2.8), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(2.6), Inches(2.85), Inches(2.6),
        items, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(5.45), Inches(12.5), Inches(0.5), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
    "情報漏えいの代償：平均被害額 約4.45百万ドル（IBM調査2023）。ブランド毀損・取引先喪失のリスクも甚大",
    sz=18, bold=True, color=GOLD)
txb(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.45),
    "3点セット：①データの棚卸し（何を持っているか把握） ②アクセス権限の整理 ③プライバシーポリシーの策定",
    sz=18, color=DKGRAY)
txb(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.45),
    "→ データを「使う」前に「守る」仕組みを整えることが、データ経営の前提条件",
    sz=18, bold=True, color=NAVY)

# ════ S19: KPI・OKRによる目標管理 ════
s = slide()
header(s, "KPI・OKRによるデータドリブン目標管理", "KPI & OKR for Data-Driven Management")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=NAVY)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.45),
    "「測定できないものは管理できない」— データで目標を設定・追跡・改善するPDCAの仕組み",
    sz=20, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.3),
    "Peter Drucker（経営の父）の言葉。データ経営の根幹となる考え方", sz=17, color=LBLUE)

rect(s, Inches(0.4), Inches(2.05), Inches(6.0), Inches(4.65), fill=LCYAN, line=BLUE, lw=Pt(1.5))
txb(s, Inches(0.6), Inches(2.15), Inches(5.6), Inches(0.5),
    "KPI（Key Performance Indicator）", sz=19, bold=True, color=NAVY)
txb(s, Inches(0.6), Inches(2.7), Inches(5.6), Inches(0.4),
    "重要業績評価指標 — 目標達成に向けた進捗を測る数値", sz=17, color=DKGRAY)
kpi_examples = [
    ("売上KPI", "月次売上・客単価・リピート率", BLUE),
    ("顧客KPI", "顧客獲得数・解約率・NPS", GREEN),
    ("業務KPI", "受注リードタイム・在庫回転率", ORANGE),
    ("財務KPI", "ROE・EBITDA・営業利益率", PURPLE),
]
for i, (cat, desc, c) in enumerate(kpi_examples):
    y = Inches(3.2) + Inches(0.75)*i
    rect(s, Inches(0.5), y, Inches(2.0), Inches(0.68), fill=c)
    txb(s, Inches(0.55), y+Inches(0.1), Inches(1.9), Inches(0.52),
        cat, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.5), y, Inches(3.7), Inches(0.68), fill=WHITE, line=c, lw=Pt(0.5))
    txb(s, Inches(2.6), y+Inches(0.1), Inches(3.5), Inches(0.52),
        desc, sz=15, color=DKGRAY)

rect(s, Inches(6.7), Inches(2.05), Inches(6.0), Inches(4.65), fill=LGREEN, line=GREEN, lw=Pt(1.5))
txb(s, Inches(6.9), Inches(2.15), Inches(5.6), Inches(0.5),
    "OKR（Objectives & Key Results）", sz=19, bold=True, color=NAVY)
txb(s, Inches(6.9), Inches(2.7), Inches(5.6), Inches(0.4),
    "目標（O）と主要な成果指標（KR）を紐付けて管理", sz=17, color=DKGRAY)
txb(s, Inches(6.9), Inches(3.2), Inches(5.6), Inches(0.45),
    "【Objective例】顧客満足度No.1になる", sz=17, bold=True, color=GREEN)
kr_items = [
    "KR1：NPS（推薦意向）を40→60に向上",
    "KR2：問い合わせ解決時間を24h→8h以内",
    "KR3：顧客満足度調査で85点以上を達成",
    "KR4：リピート購入率を60%→75%に改善",
]
txb(s, Inches(6.9), Inches(3.75), Inches(5.6), Inches(2.6),
    kr_items, sz=17, color=DKGRAY)
txb(s, Inches(6.9), Inches(6.4), Inches(5.6), Inches(0.25),
    "Google・Amazonで採用。四半期ごとに設定・評価", sz=14, color=MGRAY)

rect(s, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.45), fill=NAVY)
txb(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
    "KPIツリー：売上＝客数×客単価×購買頻度 — 分解して各指標を管理することでボトルネックを特定できる",
    sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════ S20: 次回予告（財務諸表） ════
s = slide()
header(s, "次回予告：財務諸表の読み方", "Next: Understanding Financial Statements")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=BLUE)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.5),
    "第2回：企業の「成績表」を読む — 財務3表（BS・PL・CF）の構造と経営分析への活用",
    sz=20, bold=True, color=WHITE)

statements = [
    ("貸借対照表\nBalance Sheet", NAVY, "企業の財産と負債",
     ["資産：現金・売掛金・固定資産",
      "負債：借入金・支払手形",
      "純資産：株主資本・利益剰余金",
      "「ある時点」の財政状態を表示"]),
    ("損益計算書\nP/L Statement", RED, "1年間の経営成績",
     ["売上高から費用を引いて利益を計算",
      "売上総利益・営業利益・経常利益",
      "最終的な当期純利益を確認",
      "「一定期間」の収益力を表示"]),
    ("キャッシュフロー計算書\nCash Flow", GREEN, "現金の流れ",
     ["営業CF：本業での現金収支",
      "投資CF：設備投資・資産売却",
      "財務CF：借入・返済・配当",
      "「黒字倒産」の防止に不可欠"]),
]
for i, (title, c, sub, items) in enumerate(statements):
    x = Inches(0.4) + Inches(4.2)*i
    rect(s, x, Inches(1.9), Inches(4.0), Inches(0.8), fill=c)
    txb(s, x+Inches(0.1), Inches(1.95), Inches(3.8), Inches(0.45),
        title, sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, x+Inches(0.1), Inches(2.4), Inches(3.8), Inches(0.25),
        sub, sz=13, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.7), Inches(4.0), Inches(2.6), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(2.8), Inches(3.8), Inches(2.4),
        items, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(5.45), Inches(12.5), Inches(0.5), fill=LCYAN, line=BLUE)
txb(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
    "3表のつながり：PL（純利益）→ BS（純資産に加算）→ CF（現金残高の変動）",
    sz=19, bold=True, color=NAVY)

txb(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.45),
    "■ 第2回の予習：身近な企業（トヨタ・ソフトバンク等）の有価証券報告書をEDINETで検索してみよう",
    sz=18, bold=True, color=DKGRAY)
txb(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
    "EDINET：https://disclosure.edinet-fsa.go.jp　（金融庁の電子開示システム・無料で閲覧可能）",
    sz=17, color=BLUE)

rect(s, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3), fill=NAVY)
txb(s, Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.25),
    "今日の内容を振り返り：データ経営 = データ×AI×人材×組織 の4要素を揃えて初めて実現する",
    sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

out = "C:/Users/saito/Downloads/経営分析/データ経営第1回A_v2.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
