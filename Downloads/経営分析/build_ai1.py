"""
AI概論1回_v2.pptx 生成スクリプト
標準本文24pt、全スライドをテキスト/図形で再構築 + シラバス2.0版 2スライド追加
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY    = RGBColor(0x1F, 0x38, 0x64)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x44, 0x72, 0xC4)
GREEN   = RGBColor(0x70, 0xAD, 0x47)
ORANGE  = RGBColor(0xFF, 0xC0, 0x00)
RED     = RGBColor(0xC0, 0x00, 0x00)
PURPLE  = RGBColor(0x70, 0x30, 0xA0)
DKGRAY  = RGBColor(0x26, 0x26, 0x26)
MGRAY   = RGBColor(0x59, 0x59, 0x59)
LGRAY   = RGBColor(0xF2, 0xF2, 0xF2)
LBLUE   = RGBColor(0xBD, 0xD7, 0xEE)
LCYAN   = RGBColor(0xE1, 0xF5, 0xFE)
LORANGE = RGBColor(0xFF, 0xF2, 0xCC)
LGREEN  = RGBColor(0xE2, 0xEF, 0xDA)
LRED    = RGBColor(0xFF, 0xD0, 0xD0)
TEAL    = RGBColor(0x00, 0x70, 0x6E)
LTEAL   = RGBColor(0xD0, 0xF0, 0xEE)
GOLD    = RGBColor(0xB8, 0x86, 0x00)

W = Inches(13.333); H = Inches(7.5)
FONT = "メイリオ"; STD = 24

prs = Presentation()
prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(sl, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    shp = sl.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line: shp.line.color.rgb = line; shp.line.width = lw
    else: shp.line.fill.background()
    return shp

def txb(sl, l, t, w, h, lines, sz=STD, bold=False, color=DKGRAY,
        align=PP_ALIGN.LEFT, wrap=True):
    box = sl.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    if isinstance(lines, str): lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(2)
        run = p.add_run()
        run.text = line; run.font.size = Pt(sz)
        run.font.bold = bold; run.font.color.rgb = color
        run.font.name = FONT
    return box

def header(sl, title, sub=None):
    rect(sl, 0, 0, W, Inches(0.85), fill=NAVY)
    txb(sl, Inches(0.35), Inches(0.08), W-Inches(0.7), Inches(0.72),
        title, sz=26, bold=True, color=WHITE)
    if sub:
        txb(sl, Inches(0.35), Inches(0.92), W-Inches(0.7), Inches(0.5),
            sub, sz=18, bold=True, color=BLUE)

def th(sl, y, cols, rh=Inches(0.55)):
    x = Inches(0.4)
    for txt, w in cols:
        rect(sl, x, y, w, rh, fill=NAVY, line=WHITE, lw=Pt(0.5))
        txb(sl, x+Inches(0.05), y, w-Inches(0.1), rh,
            [txt], sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w

def tr(sl, y, cells, rh=Inches(0.65), even=True):
    bg = LGRAY if even else WHITE; x = Inches(0.4)
    for txt, w, *rest in cells:
        clr = rest[0] if rest else bg
        rect(sl, x, y, w, rh, fill=clr, line=BLUE, lw=Pt(0.5))
        txb(sl, x+Inches(0.05), y, w-Inches(0.1), rh,
            [txt], sz=18, color=DKGRAY, align=PP_ALIGN.CENTER)
        x += w

# ════ S1: タイトル ════
s = slide()
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.06), fill=BLUE)
txb(s, Inches(1.5), Inches(0.7), Inches(10.3), Inches(1.4),
    "AI概論", sz=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(2.7), Inches(10.3), Inches(1.0),
    "第 1 回：AI（人工知能）とは", sz=30, color=ORANGE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6),
    "齋藤 邦彦", sz=22, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.5),
    "AIの歴史・生成AI・LLM・産業応用 — AIの全体像を一回で学ぶ", sz=18, color=LBLUE, align=PP_ALIGN.CENTER)

# ════ S2: 本日の目標 ════
s = slide()
header(s, "本日の目標")
txb(s, Inches(0.6), Inches(1.0), Inches(12.3), Inches(0.45),
    "このセッションを終了する時点で、以下のことができるようになります：",
    sz=20, color=MGRAY)
items = [
    "・ AIの定義と人間の知性との関係を説明できる",
    "・ AIの歴史（第1〜第4次ブーム）を時系列で説明できる",
    "・ ディープラーニングが第3次ブームを起こした理由を理解できる",
    "・ 生成AI・LLM・ChatGPTの概要を説明できる",
    "・ 産業別（金融・農漁業・工場・商店）のAI活用例を挙げられる",
    "・ AIを学ぶために必要な知識（統計学・Python）を理解できる",
    "・ 【New】Claude Code でAI分析・プロンプト設計ができる",
]
for i, item in enumerate(items):
    clr = LCYAN if i == len(items)-1 else LGRAY
    c = TEAL if i == len(items)-1 else DKGRAY
    b = i == len(items)-1
    rect(s, Inches(0.5), Inches(1.55)+Inches(0.7)*i, Inches(12.3), Inches(0.62),
         fill=clr, line=BLUE, lw=Pt(0.5))
    txb(s, Inches(0.65), Inches(1.6)+Inches(0.7)*i, Inches(12.0), Inches(0.58),
        item, sz=20, bold=b, color=c)

# ════ S3: AIとは ════
s = slide()
header(s, "AIとは — 人工知能の定義", "What is AI?")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(1.1), fill=NAVY)
txb(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.55),
    "AI（Artificial Intelligence）= 人工知能", sz=26, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.45),
    "人が実現するさまざまな知覚・知性・判断を人工的に再現するコンピュータ技術", sz=19, color=LBLUE)

ai_aspects = [
    ("知覚・認識", BLUE,
     ["視覚：画像・映像の認識（物体・顔・文字）",
      "聴覚：音声認識・音楽生成",
      "言語：自然言語理解・翻訳・要約"]),
    ("推論・学習", GREEN,
     ["パターン学習（機械学習）",
      "深層学習（ディープラーニング）",
      "強化学習（ゲーム・ロボット制御）"]),
    ("生成・創造", ORANGE,
     ["テキスト生成（LLM・ChatGPT）",
      "画像生成（Stable Diffusion等）",
      "動画・音声・3Dコンテンツ生成"]),
    ("判断・最適化", PURPLE,
     ["自動運転・経路最適化",
      "医療診断支援",
      "金融リスク評価・投資判断"]),
]
for i, (title, c, items) in enumerate(ai_aspects):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(2.3), Inches(3.05), Inches(0.65), fill=c)
    txb(s, x+Inches(0.1), Inches(2.35), Inches(2.85), Inches(0.6),
        title, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.95), Inches(3.05), Inches(2.3), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(3.05), Inches(2.85), Inches(2.15),
        items, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(5.45), Inches(12.5), Inches(0.6), fill=LCYAN, line=BLUE, lw=Pt(1))
txb(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.55),
    "AGI（汎用人工知能）への道：現在のAIは「特化型AI」。AGIは人間と同等の汎用知性を持つAIの第一歩",
    sz=19, color=NAVY)
txb(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.5),
    "ディープラーニング技術によってAIの精度が大幅に向上し、生成AIにより日常生活でもAIを使う時代に",
    sz=19, color=DKGRAY)
txb(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
    "→ AIは特定タスクで人間を超える能力を持つ段階に到達（画像認識・囲碁・タンパク質構造予測等）",
    sz=19, bold=True, color=NAVY)

# ════ S4: AIの歴史 ════
s = slide()
header(s, "AIの歴史：4つのブーム", "History of AI — 4 Waves")
booms = [
    ("第1次ブーム", "1950〜1960年代", "推論と探索", BLUE,
     ["チェス・数学定理証明への応用",
      "特定問題の最適解探索",
      "限界：現実の複雑な問題に対応できず",
      "→ AIの「冬の時代」（1970年代）"]),
    ("第2次ブーム", "1980〜1990年代初頭", "エキスパートシステム", GREEN,
     ["専門知識をコンピュータに実装",
      "条件式で複雑な問題を解決",
      "限界：全事例対応の難しさが露呈",
      "→ 再び「冬の時代」（1990年代）"]),
    ("第3次ブーム", "2010年代〜現在", "ディープラーニング", ORANGE,
     ["2012年画像認識コンテストで革命",
      "ビッグデータ×機械学習の実用化",
      "画像・音声・翻訳で精度が急上昇",
      "→ AI実用化の本格時代"]),
    ("第4次ブーム", "2018年〜現在", "生成AI・LLM", PURPLE,
     ["GAN・大規模言語モデルの登場",
      "ChatGPT（2022年）が世界を変えた",
      "テキスト・画像・動画を自在に生成",
      "→ AIが日常生活・業務に浸透"]),
]
for i, (boom, period, key, c, items) in enumerate(booms):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(1.05), Inches(3.05), Inches(0.65), fill=c)
    txb(s, x+Inches(0.1), Inches(1.1), Inches(2.85), Inches(0.55),
        boom, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(1.7), Inches(3.05), Inches(0.5), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(1.75), Inches(2.85), Inches(0.45),
        period, sz=16, color=MGRAY, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.2), Inches(3.05), Inches(0.55), fill=c)
    txb(s, x+Inches(0.1), Inches(2.25), Inches(2.85), Inches(0.5),
        key, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.75), Inches(3.05), Inches(3.2), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(2.85), Inches(2.85), Inches(3.0),
        items, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.55), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.5),
    "現在の位置：第4次ブーム — 生成AI×LLMが産業・社会・学術を同時に変革中",
    sz=19, bold=True, color=GOLD)

# ════ S5: ディープラーニングとは ════
s = slide()
header(s, "第3次ブームを起こしたディープラーニング", "Deep Learning Revolution")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.55), fill=NAVY)
txb(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
    "2012年 — 画像認識コンテスト（ImageNet LSVRC）でニューラルネットが圧倒的勝利",
    sz=20, bold=True, color=WHITE)

rect(s, Inches(0.4), Inches(1.7), Inches(6.0), Inches(5.0), fill=LCYAN, line=BLUE, lw=Pt(1.5))
txb(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(0.55),
    "なぜディープラーニングが強いのか？", sz=21, bold=True, color=NAVY)
reasons = [
    ("大量データ活用", "ビッグデータから自動的に特徴を学習"),
    ("多層ネットワーク", "何百・何千層のニューラル層で複雑なパターンを表現"),
    ("GPU並列計算", "従来のCPUと比べて100倍以上の処理速度"),
    ("転移学習", "一度学習したモデルを別タスクに再利用可能"),
]
for i, (r, d) in enumerate(reasons):
    y = Inches(2.45) + Inches(1.05)*i
    rect(s, Inches(0.5), y, Inches(2.0), Inches(0.85), fill=BLUE)
    txb(s, Inches(0.55), y+Inches(0.1), Inches(1.9), Inches(0.75),
        r, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.5), y, Inches(3.7), Inches(0.85), fill=WHITE, line=BLUE, lw=Pt(0.5))
    txb(s, Inches(2.6), y+Inches(0.1), Inches(3.5), Inches(0.75),
        d, sz=17, color=DKGRAY)

rect(s, Inches(6.7), Inches(1.7), Inches(6.2), Inches(5.0), fill=LGRAY, line=ORANGE, lw=Pt(1.5))
txb(s, Inches(6.9), Inches(1.8), Inches(5.8), Inches(0.55),
    "AIの進化の系譜", sz=21, bold=True, color=ORANGE)
timeline = [
    ("1990年代後半〜", "機械学習の実用化", "教師あり学習・SVMなど"),
    ("2012年", "ディープラーニング革命", "ImageNetで従来手法を大幅凌駕"),
    ("2014年", "GAN登場", "画像生成AIの原点"),
    ("2017年", "Transformer登場", "「Attention is All You Need」論文"),
    ("2018年", "BERT・GPTシリーズ", "大規模言語モデルの時代"),
    ("2022年", "ChatGPT公開", "生成AIが一般普及（史上最速で1億ユーザー）"),
    ("2024〜", "マルチモーダルAI", "テキスト・画像・音声・動画を統合"),
]
for i, (yr, event, desc) in enumerate(timeline):
    y = Inches(2.45) + Inches(0.65)*i
    rect(s, Inches(6.8), y, Inches(1.7), Inches(0.55), fill=ORANGE if yr in ["2012年","2022年"] else NAVY)
    txb(s, Inches(6.85), y+Inches(0.07), Inches(1.6), Inches(0.48),
        yr, sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(8.5), y, Inches(4.2), Inches(0.55), fill=LORANGE if yr in ["2012年","2022年"] else WHITE, line=ORANGE, lw=Pt(0.3))
    txb(s, Inches(8.6), y+Inches(0.02), Inches(4.0), Inches(0.55),
        f"{event}：{desc}", sz=14, color=DKGRAY)

# ════ S6: 生成AIとLLM ════
s = slide()
header(s, "第4次ブーム：生成AIとLLM", "Generative AI & LLM")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.95), fill=NAVY)
txb(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5),
    "生成AI = 既存データに基づいてテキスト・画像・音声・動画などを生成できるAI",
    sz=20, bold=True, color=WHITE)
txb(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.4),
    "LLM（Large Language Model）= テキストを大規模に学習し自動生成する言語モデル",
    sz=19, color=LBLUE)

llm_items = [
    ("Transformer\n（2017年）", BLUE,
     ["「Attention is All You Need」論文",
      "注意機構（アテンション）で並列処理を実現",
      "RNN・LSTMの逐次処理を超えた"]),
    ("BERT\n（2018年Google）", GREEN,
     ["Transformerのエンコーダ部分を活用",
      "文の双方向理解が可能に",
      "質問応答・文章分類で高精度"]),
    ("GPT/ChatGPT\n（OpenAI）", ORANGE,
     ["Transformerのデコーダ部分を活用",
      "入力から直接回答を生成",
      "ChatGPT3は無料で世界に普及"]),
    ("Claude\n（Anthropic）", PURPLE,
     ["Constitutional AI採用",
      "安全性・倫理重視の設計",
      "長文理解・コード生成に強み"]),
]
for i, (name, c, items) in enumerate(llm_items):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(2.1), Inches(3.05), Inches(0.9), fill=c)
    txb(s, x+Inches(0.1), Inches(2.15), Inches(2.85), Inches(0.85),
        name, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(3.0), Inches(3.05), Inches(2.4), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(3.1), Inches(2.85), Inches(2.2),
        items, sz=17, color=DKGRAY)

txb(s, Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.45),
    "ChatGPTの主な能力：", sz=20, bold=True, color=NAVY)
caps = [
    ("文章生成・要約", "レポート・メール・企画書の自動作成"),
    ("翻訳・多言語", "高精度な英語⇔日本語変換"),
    ("プログラミング", "Python・Excel・HTMLコードを生成"),
    ("英会話練習", "音声機能で英会話トレーニング"),
]
for i, (cap, desc) in enumerate(caps):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(6.1), Inches(3.05), Inches(0.55), fill=NAVY)
    txb(s, x+Inches(0.1), Inches(6.15), Inches(2.85), Inches(0.5),
        cap, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(6.65), Inches(3.05), Inches(0.65), fill=LGRAY, line=NAVY, lw=Pt(0.3))
    txb(s, x+Inches(0.08), Inches(6.7), Inches(2.9), Inches(0.6),
        desc, sz=15, color=DKGRAY, align=PP_ALIGN.CENTER)

# ════ S7: AIを学ぶために ════
s = slide()
header(s, "AIを学ぶために必要な知識", "Learning Path for AI")
txb(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.45),
    "AIを理解し活用するために必要な3つの柱：", sz=20, bold=True, color=NAVY)

pillars = [
    ("データサイエンス", BLUE,
     "データ処理・分析の技術",
     ["CSVデータの読み込み・クリーニング",
      "統計分析・可視化（グラフ・ダッシュボード）",
      "機械学習モデルの構築・評価",
      "AIに必要なデータの前処理"]),
    ("統計学・数学", GREEN,
     "AIの基礎となる学問",
     ["確率論・統計的推論",
      "線形代数（行列・ベクトル）",
      "微分・最適化（勾配降下法）",
      "文系学生も必要！概念理解が重要"]),
    ("プログラミング", ORANGE,
     "AIを実装・活用するツール",
     ["Python（NumPy・pandas・sklearn）",
      "Excel（統計関数・マクロ）",
      "ChatGPT/Claude API活用",
      "Jupyter Notebook・Google Colab"]),
]
for i, (title, c, sub, items) in enumerate(pillars):
    x = Inches(0.4) + Inches(4.2)*i
    rect(s, x, Inches(1.55), Inches(4.0), Inches(0.65), fill=c)
    txb(s, x+Inches(0.1), Inches(1.6), Inches(3.8), Inches(0.6),
        title, sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.2), Inches(4.0), Inches(0.5), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(2.25), Inches(3.8), Inches(0.45),
        sub, sz=16, color=MGRAY, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.7), Inches(4.0), Inches(2.7), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.15), Inches(2.8), Inches(3.7), Inches(2.55),
        items, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.5), fill=ORANGE)
txb(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.45),
    "本講義のロードマップ：", sz=19, bold=True, color=WHITE)
roadmap = [
    ("第1回", "AI概論", "AIの全体像・歴史・生成AI"),
    ("第2回", "データサイエンス入門", "データ経営・DX・産業応用"),
    ("第3回", "統計学基礎", "平均・分散・Σ記号・Excel"),
    ("第8-10回", "機械学習・DL", "回帰・分類・ニューラルネット"),
    ("第11回〜", "生成AI実践", "ChatGPT・Claude・プロンプト設計"),
]
th(s, Inches(6.15), [("回", Inches(1.5)), ("テーマ", Inches(3.0)), ("内容", Inches(7.7))], rh=Inches(0.48))
for i, (r, t, c) in enumerate(roadmap):
    tr(s, Inches(6.63)+Inches(0.52)*i, [(r, Inches(1.5)), (t, Inches(3.0)), (c, Inches(7.7))],
       rh=Inches(0.5), even=(i%2==0))

# ════ S8: 生成AI（画像・音声） ════
s = slide()
header(s, "生成AIの種類：画像・音声・3D・動画", "Types of Generative AI")
types = [
    ("画像生成AI", BLUE,
     ["ChatGPT（DALL-E 3）",
      "Midjourney",
      "Stable Diffusion",
      "Adobe Firefly"],
     "テキストから高品質な画像を生成\nマーケティング・デザイン・映像制作"),
    ("音声・音楽生成AI", GREEN,
     ["ElevenLabs（音声クローン）",
      "VOICEVOX（日本語TTS）",
      "Suno AI（楽曲生成）",
      "OpenAI Whisper（音声認識）"],
     "テキスト→音声変換・音声クローン\nナレーション・音楽制作・翻訳"),
    ("3D・XR・メタバース", ORANGE,
     ["Unreal Engine（ゲーム・映像）",
      "Unity（3Dゲーム・XR）",
      "Blender + ChatGPT（3Dモデル）",
      "デジタルツイン技術"],
     "AR（拡張現実）・VR（仮想現実）\nメタバース・デジタルツイン構築"),
    ("動画生成AI", PURPLE,
     ["Sora（OpenAI）",
      "Runway Gen-3",
      "Kling（快手）",
      "Lumiere（Google）"],
     "テキスト・画像から動画を生成\n広告・映像・教育コンテンツ"),
]
for i, (title, c, tools, desc) in enumerate(types):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(1.05), Inches(3.05), Inches(0.65), fill=c)
    txb(s, x+Inches(0.1), Inches(1.1), Inches(2.85), Inches(0.6),
        title, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(1.7), Inches(3.05), Inches(1.8), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(1.78), Inches(2.85), Inches(1.7),
        tools, sz=16, color=DKGRAY)
    rect(s, x, Inches(3.5), Inches(3.05), Inches(1.1), fill=c)
    txb(s, x+Inches(0.1), Inches(3.55), Inches(2.85), Inches(1.0),
        desc, sz=14, color=WHITE)

rect(s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(0.55), fill=LCYAN, line=BLUE)
txb(s, Inches(0.5), Inches(4.75), Inches(12.3), Inches(0.5),
    "マルチモーダルAI：テキスト・画像・音声・動画を統合して扱う次世代AI（GPT-4o・Gemini等）",
    sz=19, color=NAVY)

txb(s, Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.45),
    "ビジネス活用の広がり：", sz=19, bold=True, color=NAVY)
biz_uses = [
    ("コンテンツ制作", "広告・SNS・カタログを自動生成", BLUE),
    ("教育・研修", "パーソナライズ学習コンテンツ", GREEN),
    ("製品開発", "3Dモデル・プロトタイプ自動生成", ORANGE),
    ("エンタメ", "映画・ゲーム・VTuber制作", PURPLE),
]
for i, (t, d, c) in enumerate(biz_uses):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(5.95), Inches(3.05), Inches(0.55), fill=c)
    txb(s, x+Inches(0.1), Inches(6.0), Inches(2.85), Inches(0.5),
        t, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(6.5), Inches(3.05), Inches(0.8), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.08), Inches(6.55), Inches(2.9), Inches(0.7),
        d, sz=15, color=DKGRAY, align=PP_ALIGN.CENTER)

# ════ S9: 産業別AIの応用 ════
s = slide()
header(s, "産業別AIの応用", "AI Applications by Industry")
industries = [
    ("ヘルスケア", BLUE,
     ["・ 病気の早期発見・画像診断AI",
      "・ 新薬開発の候補化合物予測",
      "・ 電子カルテの自動要約",
      "・ 患者リスク層別化モデル"]),
    ("金融", GREEN,
     ["・ 信用リスク評価・不正検出",
      "・ アルゴリズム取引・高頻度取引",
      "・ 顧客セグメンテーション",
      "・ 投資戦略・ポートフォリオ最適化"]),
    ("製造・工場", ORANGE,
     ["・ スマートファクトリー（IoT×AI）",
      "・ 異常検知・故障予測・予防保全",
      "・ 工程の総合管理・最適化",
      "・ 品質検査の自動化（画像AI）"]),
    ("農漁業", PURPLE,
     ["・ ハウス農業（温度・CO2管理）",
      "・ 漁場予測（海水温画像処理）",
      "・ 農業ロボット・ドローン管理",
      "・ 需要予測による廃棄ロス削減"]),
    ("商店・流通", RED,
     ["・ 売れ筋分析（統計・可視化）",
      "・ 需要予測・在庫最適化",
      "・ 顧客行動分析・推薦システム",
      "・ 価格最適化（ダイナミックP）"]),
    ("経営・AI経営", TEAL,
     ["・ 事業計画書・報告書の自動生成",
      "・ 競合分析・SWOT自動作成",
      "・ データドリブン意思決定支援",
      "・ 中小企業診断士育成支援"]),
]
for i, (ind, c, items) in enumerate(industries):
    col = i % 3; row = i // 3
    x = Inches(0.4) + Inches(4.3)*col
    y = Inches(1.05) + Inches(2.85)*row
    rect(s, x, y, Inches(4.1), Inches(0.6), fill=c)
    txb(s, x+Inches(0.1), y+Inches(0.08), Inches(3.9), Inches(0.55),
        ind, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, y+Inches(0.6), Inches(4.1), Inches(2.1), fill=LGRAY, line=c, lw=Pt(0.8))
    txb(s, x+Inches(0.1), y+Inches(0.7), Inches(3.9), Inches(2.0),
        items, sz=17, color=DKGRAY)

# ════ S10: まとめ ════
s = slide()
header(s, "第1回 まとめ", "Lecture Summary")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.55), fill=BLUE)
txb(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
    "今日学んだキーワード", sz=20, bold=True, color=WHITE)

keywords = [
    ("AI（人工知能）", "人間の知覚・推論・判断を人工的に再現する技術"),
    ("4次ブーム", "探索→エキスパート→ディープラーニング→生成AI"),
    ("ディープラーニング", "多層ニューラルネットで大量データから自動学習する技術"),
    ("生成AI", "テキスト・画像・音声・動画を自動生成するAI"),
    ("LLM", "大規模言語モデル。Transformer・BERT・GPT・Claudeなど"),
    ("産業応用", "ヘルスケア・金融・工場・農漁業・商店×AIの具体例"),
    ("AI学習ロードマップ", "データサイエンス・統計学・Pythonの3本柱で学ぶ"),
]
for i, (kw, desc) in enumerate(keywords):
    y = Inches(1.7) + Inches(0.72)*i
    rect(s, Inches(0.4), y, Inches(3.0), Inches(0.65), fill=NAVY)
    txb(s, Inches(0.5), y+Inches(0.08), Inches(2.8), Inches(0.58),
        kw, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.4), y, Inches(9.5), Inches(0.65),
         fill=LGRAY if i%2==0 else WHITE, line=BLUE, lw=Pt(0.3))
    txb(s, Inches(3.5), y+Inches(0.08), Inches(9.3), Inches(0.58),
        desc, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.55), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
    "次回：データサイエンス入門 — データ経営・データドリブン意思決定の基礎",
    sz=20, bold=True, color=GOLD)

# ════ S11: 機械学習の種類 ════
s = slide()
header(s, "機械学習の3つのアプローチ", "Types of Machine Learning")
ml_types = [
    ("教師あり学習", BLUE, "Supervised Learning",
     ["正解ラベル付きデータで学習",
      "入力→出力のパターンを学ぶ",
      "例：スパムフィルター・画像分類",
      "例：住宅価格予測・医療診断"],
     ["回帰（Regression）",
      "分類（Classification）",
      "決定木・SVM・ランダムフォレスト"]),
    ("教師なし学習", GREEN, "Unsupervised Learning",
     ["正解なしでデータのパターンを発見",
      "データ構造の隠れた規則を見つける",
      "例：顧客セグメンテーション",
      "例：異常検知・推薦システム"],
     ["クラスタリング（k-means）",
      "次元削減（PCA・t-SNE）",
      "自己符号化器（Autoencoder）"]),
    ("強化学習", ORANGE, "Reinforcement Learning",
     ["試行錯誤で報酬を最大化",
      "エージェントが環境と相互作用",
      "例：囲碁（AlphaGo）・将棋AI",
      "例：自動運転・ロボット制御"],
     ["Q学習・方策勾配法",
      "Deep Q-Network（DQN）",
      "ChatGPT（RLHF）にも応用"]),
]
for i, (title, c, sub, items, algos) in enumerate(ml_types):
    x = Inches(0.4) + Inches(4.3)*i
    rect(s, x, Inches(1.05), Inches(4.1), Inches(0.65), fill=c)
    txb(s, x+Inches(0.1), Inches(1.1), Inches(3.9), Inches(0.35),
        title, sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, x+Inches(0.1), Inches(1.45), Inches(3.9), Inches(0.25),
        sub, sz=14, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(1.7), Inches(4.1), Inches(2.5), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(1.8), Inches(3.9), Inches(2.3),
        items, sz=17, color=DKGRAY)
    rect(s, x, Inches(4.2), Inches(4.1), Inches(1.6), fill=WHITE, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(4.3), Inches(3.9), Inches(1.4),
        algos, sz=16, color=MGRAY)

rect(s, Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.65), fill=LCYAN, line=BLUE, lw=Pt(1))
txb(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.55),
    "ディープラーニング（深層学習）は教師あり・教師なし・強化学習すべてに応用可能な汎用技術",
    sz=19, bold=True, color=NAVY)
txb(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.55),
    "→ 生成AI（ChatGPT等）は主に教師あり学習＋RLHF（人間フィードバック強化学習）で構築",
    sz=18, color=DKGRAY)

# ════ S12: ニューラルネットワーク ════
s = slide()
header(s, "ニューラルネットワークの仕組み", "How Neural Networks Work")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.75), fill=NAVY)
txb(s, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.55),
    "人間の脳の神経回路を模倣した数学モデル — 入力層・隠れ層・出力層で構成", sz=20, bold=True, color=WHITE)

layers = [
    ("入力層\nInput Layer", BLUE,
     ["特徴量（ピクセル・数値・テキスト）を受け取る",
      "各ノード＝1つの入力変数",
      "例：28×28画像→784ノード"]),
    ("隠れ層\nHidden Layer", PURPLE,
     ["重み（Weight）×入力を計算",
      "活性化関数で非線形変換",
      "多層化→ディープラーニング"]),
    ("出力層\nOutput Layer", GREEN,
     ["最終的な予測結果を出力",
      "分類：ソフトマックス関数",
      "回帰：線形出力"]),
]
for i, (title, c, items) in enumerate(layers):
    x = Inches(0.4) + Inches(4.3)*i
    rect(s, x, Inches(2.0), Inches(4.1), Inches(0.7), fill=c)
    txb(s, x+Inches(0.1), Inches(2.05), Inches(3.9), Inches(0.65),
        title, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.7), Inches(4.1), Inches(2.2), fill=LGRAY, line=c, lw=Pt(1))
    txb(s, x+Inches(0.1), Inches(2.8), Inches(3.9), Inches(2.0),
        items, sz=17, color=DKGRAY)

txb(s, Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.45),
    "■ ディープラーニングの学習プロセス", sz=20, bold=True, color=NAVY)
dl_steps = [
    ("順伝播", "入力→隠れ層→出力を計算", BLUE),
    ("損失計算", "正解と予測の誤差（Loss）", RED),
    ("逆伝播", "誤差を各層に伝達", ORANGE),
    ("重み更新", "勾配降下法で自動最適化", GREEN),
]
for i, (step, desc, c) in enumerate(dl_steps):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(5.65), Inches(3.05), Inches(0.55), fill=c)
    txb(s, x+Inches(0.1), Inches(5.7), Inches(2.85), Inches(0.5),
        step, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(6.2), Inches(3.05), Inches(0.7), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.08), Inches(6.25), Inches(2.9), Inches(0.6),
        desc, sz=16, color=DKGRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35), fill=LCYAN, line=BLUE)
txb(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
    "この繰り返し（エポック）で精度が向上 → 大量データ×GPU処理で驚異的な性能を実現",
    sz=17, bold=True, color=NAVY)

# ════ S13: 生成AIの活用事例 ════
s = slide()
header(s, "生成AIの活用事例", "Generative AI Use Cases")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=BLUE)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.5),
    "テキスト・画像・音声・動画・コードを「生成」するAI — 2022年以降に急速に普及",
    sz=20, bold=True, color=WHITE)

cases = [
    ("テキスト生成", BLUE,
     "ChatGPT / Claude / Gemini",
     ["文章作成・要約・翻訳",
      "プログラムコード生成",
      "質問応答・チャットボット",
      "報告書・企画書の自動作成"]),
    ("画像生成", PURPLE,
     "Stable Diffusion / DALL-E / Midjourney",
     ["広告デザイン・イラスト制作",
      "製品デザインのプロトタイプ",
      "医療画像の生成・拡張",
      "ゲームキャラクター・背景生成"]),
    ("音声・動画", ORANGE,
     "Suno / ElevenLabs / Runway",
     ["AI音楽・BGM自動作成",
      "テキストから動画を生成",
      "バーチャル講師・デジタル人間",
      "字幕・吹き替え自動化"]),
    ("コード生成", GREEN,
     "GitHub Copilot / Claude Code",
     ["コード補完・バグ修正",
      "テストコード自動生成",
      "APIドキュメント作成",
      "システム設計・リファクタリング"]),
]
for i, (cat, c, tool, items) in enumerate(cases):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(1.9), Inches(3.05), Inches(0.55), fill=c)
    txb(s, x+Inches(0.08), Inches(1.95), Inches(2.9), Inches(0.5),
        cat, sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.45), Inches(3.05), Inches(0.55), fill=WHITE, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.08), Inches(2.5), Inches(2.9), Inches(0.45),
        tool, sz=13, color=MGRAY, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(3.0), Inches(3.05), Inches(2.6), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(3.1), Inches(2.85), Inches(2.4),
        items, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(5.75), Inches(12.5), Inches(0.65), fill=LORANGE, line=ORANGE)
txb(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.55),
    "市場規模：生成AI市場は2023年の約400億ドルから2030年には約1.3兆ドルに成長予測（Goldman Sachs）",
    sz=19, bold=True, color=GOLD)
txb(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
    "→ あらゆる産業でAIが「道具」となり、使いこなせる人材の価値が急上昇している",
    sz=19, color=DKGRAY)

# ════ S14: AIと倫理・社会問題 ════
s = slide()
header(s, "AIと倫理・社会的課題", "AI Ethics & Social Issues")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=RED)
txb(s, Inches(0.6), Inches(1.13), Inches(12.1), Inches(0.5),
    "AIの普及とともに顕在化するリスクと倫理的課題 — 技術者・利用者として理解が必須",
    sz=20, bold=True, color=WHITE)

issues = [
    ("バイアス・差別", RED,
     ["学習データに含まれる偏見を継承",
      "人種・性別・年齢への不公平な判断",
      "採用・融資・司法でのAI差別事例",
      "→ 多様なデータと公平性検証が必要"]),
    ("フェイク・誤情報", PURPLE,
     ["ディープフェイク動画の悪用",
      "AIが生成した偽ニュース拡散",
      "なりすまし詐欺への応用",
      "→ 透明性と出典明示が重要"]),
    ("プライバシー", ORANGE,
     ["個人データの無断学習・利用",
      "顔認識による監視社会の懸念",
      "GDPR（EU）・個人情報保護法",
      "→ データ最小化と同意管理"]),
    ("雇用・経済格差", GOLD,
     ["ルーティン業務のAI代替",
      "ホワイトカラーの仕事への影響",
      "AI格差（持つ者・持たざる者）",
      "→ リスキリングと教育が鍵"]),
]
for i, (title, c, items) in enumerate(issues):
    x = Inches(0.4) + Inches(3.2)*i
    rect(s, x, Inches(1.9), Inches(3.05), Inches(0.6), fill=c)
    txb(s, x+Inches(0.08), Inches(1.95), Inches(2.9), Inches(0.55),
        title, sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.5), Inches(3.05), Inches(2.85), fill=LGRAY, line=c, lw=Pt(0.5))
    txb(s, x+Inches(0.1), Inches(2.6), Inches(2.85), Inches(2.65),
        items, sz=17, color=DKGRAY)

rect(s, Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.65), fill=LCYAN, line=BLUE)
txb(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.55),
    "AI開発原則（Google・Microsoft・Anthropic）：安全性・公平性・透明性・プライバシー保護・説明責任",
    sz=19, bold=True, color=NAVY)
txb(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
    "日本：AI戦略2022・AI安全保障・内閣府「人間中心のAI社会原則」— 国際ルール形成に参加",
    sz=18, color=DKGRAY)
txb(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45),
    "→ AIを使う際は「誰が・何のために・どんなデータを使うか」を常に意識すること",
    sz=18, bold=True, color=NAVY)

# ════ S15: 日本のAI政策と今後の動向 ════
s = slide()
header(s, "日本のAI政策と世界の動向", "Japan AI Policy & Global Trends")
txb(s, Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.45),
    "■ 日本のAI推進政策", sz=20, bold=True, color=NAVY)

policies = [
    ("AI戦略2022", BLUE, "内閣府",
     "2022年策定。デジタル社会実現に向けAI人材育成・研究開発・社会実装を推進"),
    ("Digi田甲子園", GREEN, "総務省",
     "地方自治体のDX・AI活用事例を発掘・表彰。農漁業・医療・行政への応用を支援"),
    ("AI安全保障", RED, "経産省",
     "AIチップ・基盤モデルの国産化推進。米国との連携でAI覇権に対応"),
    ("教育DX", ORANGE, "文科省",
     "大学・高校でのAIリテラシー教育を必修化。生成AI活用ガイドラインを策定"),
]
for i, (name, c, org, desc) in enumerate(policies):
    y = Inches(1.55) + Inches(0.9)*i
    rect(s, Inches(0.4), y, Inches(2.5), Inches(0.8), fill=c)
    txb(s, Inches(0.5), y+Inches(0.08), Inches(2.3), Inches(0.35),
        name, sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, Inches(0.5), y+Inches(0.45), Inches(2.3), Inches(0.3),
        org, sz=13, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.9), y, Inches(10.0), Inches(0.8), fill=LGRAY if i%2==0 else WHITE, line=c, lw=Pt(0.5))
    txb(s, Inches(3.0), y+Inches(0.15), Inches(9.8), Inches(0.6),
        desc, sz=18, color=DKGRAY)

txb(s, Inches(0.4), Inches(5.2), Inches(12.5), Inches(0.45),
    "■ 世界のAIトレンド 2025-2030", sz=20, bold=True, color=NAVY)
trends = [
    ("AGIへの競争", "OpenAI・Google・Anthropicが汎用人工知能の実現を目指して開発競争", BLUE),
    ("マルチモーダルAI", "テキスト・画像・音声・動画を横断して処理するAIモデルの実用化", PURPLE),
    ("AIエージェント", "自律的にタスクを実行するAIエージェントが企業業務を変革", ORANGE),
    ("省エネAI", "大規模モデルの計算コスト・電力消費削減が次の技術課題に", GREEN),
]
for i, (title, desc, c) in enumerate(trends):
    y = Inches(5.75) + Inches(0.42)*i
    rect(s, Inches(0.4), y, Inches(2.5), Inches(0.38), fill=c)
    txb(s, Inches(0.5), y+Inches(0.04), Inches(2.3), Inches(0.34),
        title, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, Inches(3.0), y+Inches(0.04), Inches(9.8), Inches(0.34),
        desc, sz=16, color=DKGRAY)

# ════ S16: 課題・次回予告 ════
s = slide()
header(s, "課題・次回予告", "Assignment & Next Lecture")
rect(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.65), fill=ORANGE)
txb(s, Inches(0.5), Inches(1.13), Inches(12.3), Inches(0.55),
    "【今週の課題】AIを活用した体験レポート（初回は提出なし・体験してみよう）",
    sz=21, bold=True, color=WHITE)

rect(s, Inches(0.4), Inches(1.85), Inches(6.0), Inches(4.95), fill=LGRAY, line=BLUE, lw=Pt(1))
txb(s, Inches(0.5), Inches(1.95), Inches(5.8), Inches(0.5),
    "■ 課題の内容", sz=20, bold=True, color=NAVY)
tasks = [
    ("TASK 1", "ChatGPTまたはClaudeに自由な質問をする（3回以上）", BLUE),
    ("TASK 2", "画像生成AI（DALL-E等）で画像を1枚生成してみる", GREEN),
    ("TASK 3", "AIの回答が正しいか自分で調べて確認する", ORANGE),
    ("TASK 4", "AIを使った感想を200字でまとめる", PURPLE),
    ("TASK 5", "「AIに向いている仕事・向かない仕事」を考える", RED),
]
for i, (t, d, c) in enumerate(tasks):
    y = Inches(2.55) + Inches(0.7)*i
    rect(s, Inches(0.5), y, Inches(1.5), Inches(0.62), fill=c)
    txb(s, Inches(0.55), y+Inches(0.1), Inches(1.4), Inches(0.52),
        t, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(2.0), y, Inches(4.2), Inches(0.62), fill=WHITE, line=c, lw=Pt(0.5))
    txb(s, Inches(2.1), y+Inches(0.1), Inches(4.0), Inches(0.52),
        d, sz=16, color=DKGRAY)

rect(s, Inches(6.6), Inches(1.85), Inches(6.3), Inches(4.95), fill=LCYAN, line=BLUE, lw=Pt(1))
txb(s, Inches(6.7), Inches(1.95), Inches(6.1), Inches(0.5),
    "■ 次回（第2回）の内容", sz=20, bold=True, color=NAVY)
next_items = [
    "📘 データサイエンスの基礎",
    "　　 — データの種類・構造・収集方法",
    "📊 データの前処理と可視化",
    "　　 — 欠損値・外れ値・グラフの読み方",
    "🐍 Pythonでデータを扱う",
    "　　 — pandas基礎・DataFrame操作",
    "🔑 キーワード：",
    "　　 量的データ・質的データ・EDA",
]
txb(s, Inches(6.7), Inches(2.55), Inches(6.1), Inches(3.6),
    next_items, sz=18, color=DKGRAY)

rect(s, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.45), fill=NAVY)
txb(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
    "次回までにChatGPTまたはClaudeのアカウントを作成しておくこと",
    sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

out = "C:/Users/saito/Downloads/経営分析/AI概論1回_v2.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
